from __future__ import annotations

import io
import tempfile
import unittest
from contextlib import redirect_stdout
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation

from tests._util import TEMPLATE

from slidemaker import SlideBuilder


class SlideBuilderTests(unittest.TestCase):
    def test_constructor_and_add_style_validation(self) -> None:
        builder = SlideBuilder(
            TEMPLATE,
            style={".slide": {"font-size": 22}, "dense": {"spacing": 8}},
            template_default_page=4,
        )
        self.assertEqual(builder._template_default_page, 4)
        self.assertEqual(builder._styles[".slide"]["font-size"], 22)
        self.assertEqual(builder._styles["dense"]["spacing"], 8)

        with self.assertRaisesRegex(TypeError, "style name keys must be strings"):
            builder.add_style({1: {"font-size": 20}})  # type: ignore[arg-type]

        with self.assertRaisesRegex(TypeError, "must map to a dictionary"):
            builder.add_style({"bad": "value"})  # type: ignore[arg-type]

    def test_resolve_styles_variants(self) -> None:
        builder = SlideBuilder(TEMPLATE)
        builder.add_style(
            {
                ".slide": {"font-size": 24},
                ".code": {"line-numbers": True},
                ".table": {"padding": "4pt"},
                "dense": {"spacing": 8},
                "caps": {"uppercase": True},
            }
        )

        styles = builder._resolve_styles(None)
        self.assertEqual(styles[".slide"]["font-size"], 24)
        self.assertTrue(styles[".code"]["line-numbers"])

        named = builder._resolve_styles("dense")
        self.assertEqual(named[".slide"]["spacing"], 8)

        nested = builder._resolve_styles(
            {
                "use": ["dense", "caps", "missing"],
                ".slide": {"font-size": 20},
                ".code": {"font-size": 18},
                ".table-header": {"bold": True},
                "#title": {"font-color": "#FFFFFF"},
                "italic": True,
            }
        )
        self.assertTrue(nested[".slide"]["uppercase"])
        self.assertEqual(nested[".slide"]["font-size"], 20)
        self.assertEqual(nested[".slide"]["italic"], True)
        self.assertEqual(nested[".code"]["font-size"], 18)
        self.assertTrue(nested[".table-header"]["bold"])
        self.assertEqual(nested["#title"]["font-color"], "#FFFFFF")

        flat = builder._resolve_styles({"font-size": 30, "bold": True})
        self.assertEqual(flat[".slide"]["font-size"], 30)
        self.assertTrue(flat[".slide"]["bold"])

        with self.assertRaisesRegex(KeyError, "unknown style name"):
            builder._resolve_styles("unknown")

        with self.assertRaisesRegex(
            TypeError, "style must be None, a style name, or a dict"
        ):
            builder._resolve_styles(123)  # type: ignore[arg-type]

        with self.assertRaisesRegex(KeyError, "unknown style name"):
            builder._resolve_styles({"use": "missing", ".slide": {"font-size": 10}})

    def test_add_slide_dispatches_helpers(self) -> None:
        builder = SlideBuilder(TEMPLATE)
        fake_slide = object()
        style = {
            ".slide": {"font-size": 24},
            ".code": {"line-numbers": True},
            ".table": {"padding": "6pt"},
            ".table-header": {"bold": True},
            ".table-cell": {"italic": True},
        }

        with (
            patch.object(
                builder, "_resolve_styles", return_value=style
            ) as resolve_styles,
            patch("slidemaker.cli.clone_slide", return_value=fake_slide) as clone_slide,
            patch("slidemaker.cli.replace_placeholders") as replace_placeholders,
            patch(
                "slidemaker.cli.remove_generated_content_placeholders"
            ) as remove_generated,
            patch("slidemaker.cli.layout_content_shapes") as layout_content_shapes,
            patch("slidemaker.cli.set_notes") as set_notes,
        ):
            builder.add_slide(
                content={"title": "Hello"},
                items=["One"],
                code="print('x')",
                table={"columns": ["A"], "rows": [["1"]]},
                flow_boxes=[{"label": "X"}],
                callout="Done",
                notes="speaker note",
                style="dense",
                template_page=2,
            )

        resolve_styles.assert_called_once_with("dense")
        clone_slide.assert_called_once_with(builder._prs, 1)
        replace_placeholders.assert_called_once_with(
            fake_slide, {"title": "Hello"}, styles=style
        )
        remove_generated.assert_called_once_with(fake_slide)
        layout_content_shapes.assert_called_once()
        _, kwargs = layout_content_shapes.call_args
        self.assertEqual(kwargs["slide_style"], style[".slide"])
        self.assertEqual(kwargs["code_style"], style[".code"])
        self.assertEqual(kwargs["table_style"], style[".table"])
        self.assertEqual(kwargs["table_header_style"], style[".table-header"])
        self.assertEqual(kwargs["table_cell_style"], style[".table-cell"])
        set_notes.assert_called_once_with(fake_slide, "speaker note")
        self.assertEqual(builder._slide_count, 1)

    def test_add_slide_skips_optional_paths(self) -> None:
        builder = SlideBuilder(TEMPLATE)
        fake_slide = object()

        with (
            patch("slidemaker.cli.clone_slide", return_value=fake_slide),
            patch("slidemaker.cli.replace_placeholders") as replace_placeholders,
            patch(
                "slidemaker.cli.remove_generated_content_placeholders"
            ) as remove_generated,
            patch("slidemaker.cli.layout_content_shapes") as layout_content_shapes,
            patch("slidemaker.cli.set_notes") as set_notes,
        ):
            builder.add_slide()

        replace_placeholders.assert_not_called()
        remove_generated.assert_not_called()
        layout_content_shapes.assert_not_called()
        set_notes.assert_not_called()

    def test_save_writes_only_generated_slides(self) -> None:
        builder = SlideBuilder(TEMPLATE, template_default_page=4)
        builder.add_slide(content={"title": "Saved Slide"}, items=["A", "B"])

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / "deck.pptx"
            stream = io.StringIO()
            with redirect_stdout(stream):
                builder.save(str(out_path))

            self.assertTrue(out_path.exists())
            saved = Presentation(str(out_path))
            self.assertEqual(len(saved.slides), 1)

        self.assertIn("Saved", stream.getvalue())


if __name__ == "__main__":
    unittest.main()
