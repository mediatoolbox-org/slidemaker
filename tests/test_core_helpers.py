from __future__ import annotations

import unittest
from types import SimpleNamespace

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from tests._util import rgb_hex

from slidemaker import core


class DummyRPr(dict):
    def set(self, key: str, value: str) -> None:
        self[key] = value


class DummyRun:
    def __init__(self) -> None:
        self._r = SimpleNamespace(get_or_add_rPr=lambda: self._r_pr)
        self._r_pr = DummyRPr()


class DummyFill:
    def __init__(self) -> None:
        self.mode: str | None = None
        self.fore_color = SimpleNamespace(rgb=None)

    def background(self) -> None:
        self.mode = "background"

    def solid(self) -> None:
        self.mode = "solid"


class DummyTarget:
    def __init__(self) -> None:
        self.fill = DummyFill()


class CoreHelperTests(unittest.TestCase):
    def test_style_normalization_and_merge(self) -> None:
        normalized = core._normalize_style(
            {"Font_Size": 24, " font_color ": "#FFFFFF", 1: "ignored"}
        )
        self.assertEqual(
            normalized,
            {"font-size": 24, "font-color": "#FFFFFF"},
        )
        self.assertEqual(core._normalize_style(None), {})

        merged = core._merge_style({"font-size": 20}, {"Font_Size": 24, "bold": True})
        self.assertEqual(merged, {"font-size": 24, "bold": True})
        self.assertEqual(core._merge_style(None, None), {})

    def test_color_and_length_parsers(self) -> None:
        default = RGBColor(9, 9, 9)
        self.assertEqual(rgb_hex(core._as_rgb_color("#010203")), "010203")
        self.assertEqual(rgb_hex(core._as_rgb_color("010203")), "010203")
        self.assertEqual(rgb_hex(core._as_rgb_color((1, 2, 3))), "010203")
        self.assertEqual(rgb_hex(core._as_rgb_color([1, 2, 3])), "010203")
        self.assertIs(core._as_rgb_color(default), default)
        self.assertEqual(core._as_rgb_color("bad", default), default)
        self.assertEqual(core._as_rgb_color((999, 1, 1), default), default)
        self.assertEqual(core._as_rgb_color(("x", 1, 1), default), default)

        self.assertEqual(core._as_pt(12).pt, 12)
        self.assertEqual(core._as_pt("13pt").pt, 13)
        self.assertEqual(core._as_pt("14").pt, 14)
        self.assertEqual(core._as_pt("bad", Pt(9)), Pt(9))
        self.assertEqual(core._as_pt(None, Pt(8)), Pt(8))

        self.assertEqual(core._as_length(7), 7)
        self.assertEqual(core._as_length(1.5), Inches(1.5))
        self.assertEqual(core._as_length("2in"), Inches(2))
        self.assertEqual(core._as_length("18pt"), Pt(18))
        self.assertEqual(core._as_length("120emu"), 120)
        self.assertEqual(core._as_length("3.5"), Inches(3.5))
        self.assertEqual(core._as_length(True, 99), 99)
        self.assertEqual(core._as_length("bad", 77), 77)
        self.assertEqual(core._as_length(SimpleNamespace(emu="bad"), 66), 66)
        self.assertEqual(core._as_length(SimpleNamespace(emu=1200)), 1200)

    def test_spacing_bool_alignment_and_uppercase_helpers(self) -> None:
        self.assertEqual(core._font_size_pt(Pt(11)), 11.0)
        self.assertIsNone(core._font_size_pt(None))
        self.assertIsNone(core._font_size_pt(SimpleNamespace(pt="bad")))

        self.assertEqual(core._resolve_line_spacing(1.5, Pt(20)).pt, 30)
        self.assertEqual(core._resolve_line_spacing("36pt", Pt(20)).pt, 36)
        self.assertEqual(core._resolve_line_spacing("2x", Pt(20)).pt, 40)
        self.assertEqual(core._resolve_line_spacing("150%", Pt(20)).pt, 30)
        self.assertEqual(core._resolve_line_spacing("1.2", Pt(20)).pt, 24)
        self.assertEqual(
            core._resolve_line_spacing("oops", Pt(20), "fallback"), "fallback"
        )
        self.assertEqual(core._resolve_line_spacing(2, object()), 2.0)

        self.assertEqual(core._resolve_letter_spacing(20, Pt(10)), 20)
        self.assertEqual(core._resolve_letter_spacing("0.9pt", Pt(10)), 90)
        self.assertEqual(core._resolve_letter_spacing("10", Pt(20)), 20)
        self.assertEqual(core._resolve_letter_spacing("oops", Pt(20), 7), 7)
        self.assertEqual(core._resolve_letter_spacing(3, object()), 3)

        self.assertTrue(core._as_bool(True))
        self.assertTrue(core._as_bool(1))
        self.assertFalse(core._as_bool("off"))
        self.assertEqual(core._as_bool("maybe", True), True)

        self.assertTrue(core._resolve_uppercase({"text-transform": "uppercase"}))
        self.assertFalse(core._resolve_uppercase({"text-transform": "normal"}, True))
        self.assertTrue(core._resolve_uppercase({"uppercase": "yes"}))
        self.assertEqual(core._apply_uppercase("abc", True), "ABC")
        self.assertEqual(core._apply_uppercase("abc", False), "abc")

        self.assertEqual(core._as_alignment("center"), PP_ALIGN.CENTER)
        self.assertEqual(core._as_alignment(PP_ALIGN.RIGHT), PP_ALIGN.RIGHT)
        self.assertEqual(core._as_alignment("unknown", PP_ALIGN.LEFT), PP_ALIGN.LEFT)

    def test_code_numbering_bold_segments_padding_and_fill(self) -> None:
        numbered = core._with_code_line_numbers("a\nb")
        self.assertIn(" 1  a", numbered)
        self.assertEqual(core._with_code_line_numbers(" 1  a\n 2  b"), " 1  a\n 2  b")

        self.assertEqual(
            core._markdown_bold_segments("before **BOLD** after"),
            [("before ", False), ("BOLD", True), (" after", False)],
        )
        self.assertEqual(core._markdown_bold_segments("plain"), [("plain", False)])
        self.assertEqual(core._markdown_bold_segments("****"), [])

        padding = core._resolve_padding(
            {
                "padding": "4pt",
                "padding-left": "6pt",
                "padding-bottom": "8pt",
            },
            Pt(2),
        )
        self.assertEqual(padding[0], Pt(6))
        self.assertEqual(padding[1], Pt(4))
        self.assertEqual(padding[2], Pt(4))
        self.assertEqual(padding[3], Pt(8))

        tf = SimpleNamespace(
            margin_left=0, margin_top=0, margin_right=0, margin_bottom=0
        )
        core._apply_text_frame_padding(tf, {"padding": "5pt"})
        self.assertEqual(tf.margin_left, Pt(5))
        self.assertEqual(tf.margin_bottom, Pt(5))

        run = DummyRun()
        core._apply_run_letter_spacing(run, None)
        self.assertEqual(run._r_pr, {})
        core._apply_run_letter_spacing(run, 123)
        self.assertEqual(run._r_pr["spc"], "123")

        target = DummyTarget()
        core._set_fill_color(target, None)
        self.assertEqual(target.fill.mode, "background")
        core._set_fill_color(target, RGBColor(1, 2, 3))
        self.assertEqual(target.fill.mode, "solid")
        self.assertEqual(rgb_hex(target.fill.fore_color.rgb), "010203")

    def test_dummy_shape_lookup_helpers(self) -> None:
        child = SimpleNamespace(name="Child", has_text_frame=True, shape_type=17)
        non_text_child = SimpleNamespace(
            name="Other", has_text_frame=False, shape_type=1
        )
        group = SimpleNamespace(
            name="Group", shape_type=6, shapes=[non_text_child, child]
        )
        textbox = SimpleNamespace(name="Standalone", has_text_frame=True, shape_type=17)
        other = SimpleNamespace(name="Shape", has_text_frame=False, shape_type=1)
        slide = SimpleNamespace(shapes=[group, textbox, other])

        self.assertIs(core.find_group_textbox(slide, "Group"), child)
        self.assertIsNone(core.find_group_textbox(slide, "Missing"))
        self.assertIs(core.find_textbox_by_name(slide, "Standalone"), textbox)
        self.assertIsNone(core.find_textbox_by_name(slide, "Missing"))

        iterated = core._iter_text_shapes(slide)
        self.assertEqual(iterated, [child, textbox])

    def test_table_spec_and_split_content_height(self) -> None:
        self.assertEqual(
            core._table_spec_value({"row_heights": [1]}, "row-heights"), [1]
        )
        self.assertEqual(
            core._table_spec_value({"row-heights": [2]}, "row_heights"), [2]
        )
        self.assertEqual(core._table_spec_value({}, "missing", default="x"), "x")

        first, second = core._split_content_height(100, 0.4, 10, 20, 30)
        self.assertEqual(first + second + 10, 100)
        self.assertGreaterEqual(first, 20)
        self.assertGreaterEqual(second, 30)

        first, second = core._split_content_height(40, 0.9, 10, 30, 20)
        self.assertEqual(first + second + 10, 40)
        self.assertGreaterEqual(first, 1)
        self.assertGreaterEqual(second, 1)


if __name__ == "__main__":
    unittest.main()
