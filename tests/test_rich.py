from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tests._util import SAMPLE_IMAGE, TEMPLATE, slide_texts

from slidemaker import SlideBuilder


class RichSlideBuilderTests(unittest.TestCase):
    def test_build_deck_with_flow_code_and_table_content(self) -> None:
        sb = SlideBuilder(TEMPLATE, template_default_page=4)

        sb.add_slide(
            content={"title": "The ETL Paradigm"},
            flow_boxes=[
                {
                    "label": "EXTRACT",
                    "desc": "Retrieve raw data\nfrom a source",
                    "style": {"fill-color": "#2E86AB"},
                },
                {
                    "label": "TRANSFORM",
                    "desc": "Clean, filter, enrich\nor reshape",
                    "style": {"fill-color": "#48A99A"},
                },
                {
                    "label": "LOAD",
                    "desc": "Write results to\na destination",
                    "style": {"fill-color": "#E86F51"},
                },
            ],
            callout="Separating stages makes each one independently testable and swappable",
            notes="ETL stands for Extract Transform Load.",
        )

        sb.add_slide(
            content={"title": "Extract: Date-Range Queries in MongoDB"},
            items=[
                "Convert date string to Timestamp with pd.to_datetime",
                "Compute end of day with pd.DateOffset(days=1)",
                "Use $gte and $lt for a half-open interval",
                "Combine date filter with field filter",
            ],
            code="""start = pd.to_datetime("2022-05-02",
                       format="%Y-%m-%d")
end = start + pd.DateOffset(days=1)
query = {
    "createdAt": {"$gte": start, "$lt": end},
    "admissionsQuiz": "incomplete",
}
results = list(collection.find(query))""",
            notes="The extract step.",
        )

        sb.add_slide(
            content={"title": "MongoDB Field Reference"},
            table={
                "columns": ["Field", "Type", "Notes"],
                "rows": [
                    ["_id", "ObjectId", "Primary key"],
                    ["createdAt", "datetime", "Stored in UTC"],
                    ["admissionsQuiz", "string", "Completion status"],
                ],
                "column_widths": [2.3, 2.0, 5.7],
                "banded_rows": True,
            },
            notes="Table layout smoke test.",
        )

        sb.add_slide(
            content={"title": "Items And Table"},
            items=[
                "Key",
                "Document",
            ],
            table={
                "columns": ["Field", "Meaning"],
                "rows": [["_id", "Primary key"], ["createdAt", "Timestamp"]],
            },
        )

        sb.add_slide(
            content={"title": "Code And Table"},
            code='doc = {"_id": 1}',
            table={
                "columns": ["Field", "Value"],
                "rows": [["_id", "1"]],
            },
        )

        sb.add_slide(
            content={"title": "Model Chart"},
            items=["Validation AUC improved", "Class separation is clearer"],
            image={
                "path": SAMPLE_IMAGE,
                "caption": "Validation confusion matrix",
            },
        )

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / "rich-smoke.pptx"
            sb.save(str(out_path))
            prs = Presentation(str(out_path))

        self.assertEqual(len(prs.slides), 6)

        flow_text = "\n".join(slide_texts(prs.slides[0]))
        self.assertIn("EXTRACT", flow_text)
        self.assertIn(
            "Separating stages makes each one independently testable", flow_text
        )

        code_text = "\n".join(slide_texts(prs.slides[1]))
        self.assertIn("start = pd.to_datetime", code_text)
        self.assertIn("Combine date filter with field filter", code_text)

        table_frame = next(shape for shape in prs.slides[2].shapes if shape.has_table)
        self.assertEqual(table_frame.table.cell(0, 0).text, "Field")
        self.assertEqual(table_frame.table.cell(1, 0).text, "_id")
        self.assertEqual(table_frame.table.cell(2, 2).text, "Stored in UTC")

        items_table_frame = next(
            shape for shape in prs.slides[3].shapes if shape.has_table
        )
        self.assertEqual(items_table_frame.table.cell(1, 1).text, "Primary key")

        code_table_text = "\n".join(slide_texts(prs.slides[4]))
        self.assertIn('doc = {"_id": 1}', code_table_text)
        code_table_frame = next(
            shape for shape in prs.slides[4].shapes if shape.has_table
        )
        self.assertEqual(code_table_frame.table.cell(1, 1).text, "1")

        image_slide_text = "\n".join(slide_texts(prs.slides[5]))
        self.assertIn("Validation AUC improved", image_slide_text)
        self.assertIn("Validation confusion matrix", image_slide_text)
        self.assertTrue(
            any(
                shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                for shape in prs.slides[5].shapes
            )
        )


if __name__ == "__main__":
    unittest.main()
