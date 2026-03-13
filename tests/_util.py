from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

TEMPLATE = ROOT / "tests" / "data" / "template.pptx"


def new_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            texts.append(shape.text_frame.text)
    return texts


def rgb_hex(value) -> str | None:
    if value is None:
        return None
    return str(value).upper()
