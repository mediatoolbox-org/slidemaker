from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from tests._util import TEMPLATE, slide_texts

from slidemaker import SlideBuilder


class SlideBuilderSmokeTests(unittest.TestCase):
    def test_build_deck_with_bullets_and_notes(self) -> None:
        sb = SlideBuilder(TEMPLATE, template_default_page=4)

        sb.add_slide(
            content={"title": "The ETL Paradigm"},
            items=[
                "Extract: Retrieve raw data from a source",
                "Transform: Clean, filter, enrich or reshape",
                "Load: Write results to a destination",
            ],
            notes="ETL stands for Extract Transform Load.",
        )

        sb.add_slide(
            content={"title": "Extract: Date-Range Queries"},
            items=[
                "Convert date string to Timestamp",
                "Compute end of day with DateOffset",
                "Use $gte and $lt for half-open interval",
            ],
            notes="The extract step.",
        )

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / "build-smoke.pptx"
            sb.save(str(out_path))
            prs = Presentation(str(out_path))

        self.assertEqual(len(prs.slides), 2)
        first_text = "\n".join(slide_texts(prs.slides[0]))
        second_text = "\n".join(slide_texts(prs.slides[1]))
        self.assertIn("Content Title", first_text)
        self.assertIn("Extract: Retrieve raw data from a source", first_text)
        self.assertIn("Content Title", second_text)
        self.assertIn("Use $gte and $lt for half-open interval", second_text)
        self.assertEqual(
            prs.slides[0].notes_slide.notes_text_frame.text,
            "ETL stands for Extract Transform Load.",
        )
        self.assertEqual(
            prs.slides[1].notes_slide.notes_text_frame.text,
            "The extract step.",
        )


if __name__ == "__main__":
    unittest.main()
