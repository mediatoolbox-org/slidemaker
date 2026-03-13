from __future__ import annotations

import unittest
from pathlib import Path
from tempfile import TemporaryDirectory
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from tests._util import SAMPLE_IMAGE, TEMPLATE, new_slide, rgb_hex, slide_texts

from slidemaker import core


class CorePptxTests(unittest.TestCase):
    def test_set_textbox_text_applies_style_and_handles_empty(self) -> None:
        _, slide = new_slide()
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1.5))

        core.set_textbox_text(
            shape,
            "hello",
            style={
                "font-size": 28,
                "font-color": "#112233",
                "font-name": "Aptos",
                "bold": True,
                "italic": True,
                "align": "center",
                "padding": "4pt",
                "letter-spacing": 30,
                "line-height": "1.5",
                "uppercase": True,
            },
        )

        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        self.assertEqual(shape.text, "HELLO")
        self.assertEqual(run.font.size.pt, 28)
        self.assertEqual(rgb_hex(run.font.color.rgb), "112233")
        self.assertEqual(run.font.name, "Aptos")
        self.assertTrue(run.font.bold)
        self.assertTrue(run.font.italic)
        self.assertEqual(paragraph.alignment, PP_ALIGN.CENTER)
        self.assertEqual(shape.text_frame.margin_left, Pt(4))
        self.assertEqual(run._r.get_or_add_rPr().get("spc"), "84")

        core.set_textbox_text(shape, "")
        self.assertEqual(shape.text, "")

    def test_add_textbox_and_bullet_list_create_styled_text(self) -> None:
        _, slide = new_slide()

        textbox = core.add_textbox(
            slide,
            Inches(0.5),
            Inches(0.5),
            Inches(3),
            Inches(1),
            "MongoDB",
            style={
                "uppercase": True,
                "font-name": "Aptos",
                "font-size": 26,
                "italic": True,
                "alignment": "center",
            },
        )
        paragraph = textbox.text_frame.paragraphs[0]
        self.assertEqual(textbox.text, "MONGODB")
        self.assertEqual(paragraph.font.size.pt, 26)
        self.assertEqual(paragraph.font.name, "Aptos")
        self.assertTrue(paragraph.font.italic)
        self.assertEqual(paragraph.alignment, PP_ALIGN.CENTER)

        bullets = core.add_bullet_list(
            slide,
            Inches(1),
            Inches(2),
            Inches(5),
            Inches(3),
            ["Before **Bold** after", "Second item"],
            style={
                "bullet-char": "-",
                "font-size": 22,
                "italic": True,
                "alignment": "right",
                "space-before": "4pt",
                "space-after": "8pt",
                "letter-spacing": "0.5pt",
            },
        )
        first_paragraph = bullets.text_frame.paragraphs[0]
        first_text = "".join(run.text for run in first_paragraph.runs)
        self.assertEqual(first_text, "Before Bold after")
        self.assertFalse(first_paragraph.runs[0].font.bold)
        self.assertTrue(first_paragraph.runs[1].font.bold)
        self.assertTrue(first_paragraph.runs[1].font.italic)
        self.assertEqual(first_paragraph.alignment, PP_ALIGN.RIGHT)
        self.assertEqual(first_paragraph.space_before.pt, 4)
        self.assertEqual(first_paragraph.space_after.pt, 8)
        self.assertIn('char="-"', first_paragraph._p.xml)
        self.assertIn('marL="', first_paragraph._p.xml)
        self.assertIn('indent="', first_paragraph._p.xml)

    def test_add_markdown_textbox_supports_paragraphs_and_inline_styles(self) -> None:
        _, slide = new_slide()

        textbox = core.add_markdown_textbox(
            slide,
            Inches(0.5),
            Inches(0.5),
            Inches(6),
            Inches(3),
            "# Heading\n\nParagraph with **bold**, *italic*, and `code`.\n  - Bullet item",
            style={
                "font-size": 22,
                "font-color": "#112233",
                "padding": "4pt",
            },
        )

        paragraphs = textbox.text_frame.paragraphs
        self.assertEqual(len(paragraphs), 3)
        self.assertEqual(paragraphs[0].text, "Heading")
        self.assertGreater(paragraphs[0].runs[0].font.size.pt, 22)
        self.assertTrue(paragraphs[0].runs[0].font.bold)
        self.assertEqual(textbox.text_frame.margin_left, Pt(4))

        paragraph_text = "".join(run.text for run in paragraphs[1].runs)
        self.assertEqual(paragraph_text, "Paragraph with bold, italic, and code.")
        bold_run = next(run for run in paragraphs[1].runs if run.text == "bold")
        italic_run = next(run for run in paragraphs[1].runs if run.text == "italic")
        code_run = next(run for run in paragraphs[1].runs if run.text == "code")
        self.assertTrue(bold_run.font.bold)
        self.assertTrue(italic_run.font.italic)
        self.assertEqual(code_run.font.name, core.CODE_FONT)
        self.assertEqual(rgb_hex(code_run.font.color.rgb), "112233")

        self.assertEqual(paragraphs[2].text, "Bullet item")
        self.assertIn("buChar", paragraphs[2]._p.xml)
        self.assertIn('lvl="1"', paragraphs[2]._p.xml)
        self.assertIn('marL="', paragraphs[2]._p.xml)
        self.assertIn('indent="', paragraphs[2]._p.xml)

    def test_add_markdown_textbox_serializes_nested_real_bullets(self) -> None:
        prs, slide = new_slide()

        textbox = core.add_markdown_textbox(
            slide,
            Inches(0.5),
            Inches(0.5),
            Inches(6),
            Inches(3),
            "- Parent\n  - Child\n    - Grandchild",
            style={"bullet-char": "-"},
        )

        paragraphs = textbox.text_frame.paragraphs
        self.assertEqual(
            [paragraph.text for paragraph in paragraphs],
            ["Parent", "Child", "Grandchild"],
        )
        self.assertIn('char="-"', paragraphs[0]._p.xml)
        self.assertIn('lvl="0"', paragraphs[0]._p.xml)
        self.assertIn('lvl="1"', paragraphs[1]._p.xml)
        self.assertIn('lvl="2"', paragraphs[2]._p.xml)

        with TemporaryDirectory() as tmpdir:
            out_path = Path(tmpdir) / "markdown-bullets.pptx"
            prs.save(out_path)
            with ZipFile(out_path) as pptx_zip:
                slide_xml = pptx_zip.read("ppt/slides/slide1.xml").decode("utf-8")

        self.assertGreaterEqual(slide_xml.count('<a:buChar char="-"/>'), 3)
        self.assertIn('lvl="0"', slide_xml)
        self.assertIn('lvl="1"', slide_xml)
        self.assertIn('lvl="2"', slide_xml)
        self.assertIn('marL="304800"', slide_xml)
        self.assertIn('marL="533400"', slide_xml)
        self.assertIn('marL="762000"', slide_xml)
        self.assertIn('indent="-228600"', slide_xml)

    def test_add_shape_rect_and_code_block(self) -> None:
        _, slide = new_slide()

        rect = core.add_shape_rect(
            slide,
            Inches(0.5),
            Inches(0.5),
            Inches(2),
            Inches(1),
            style={
                "fill-color": "#123456",
                "line-color": "#654321",
                "line-width": "2pt",
            },
        )
        self.assertEqual(rgb_hex(rect.fill.fore_color.rgb), "123456")
        self.assertEqual(rgb_hex(rect.line.color.rgb), "654321")
        self.assertEqual(rect.line.width.pt, 2)

        core.add_shape_rect(
            slide,
            Inches(3),
            Inches(0.5),
            Inches(2),
            Inches(1),
        )

        core.add_code_block(
            slide,
            Inches(0.5),
            Inches(2),
            Inches(5),
            Inches(2.5),
            "a = 1\nb = 2",
            style={
                "bg-color": "#193952",
                "font-color": "#FFFFFF",
                "line-numbers": True,
                "uppercase": True,
                "padding": "6pt",
                "letter-spacing": 20,
            },
        )
        code_background = slide.shapes[2]
        code_text = slide.shapes[3]
        self.assertEqual(rgb_hex(code_background.fill.fore_color.rgb), "193952")
        self.assertIn(" 1  A = 1", code_text.text)
        self.assertIn(" 2  B = 2", code_text.text)

    def test_add_image_supports_contain_and_stretch(self) -> None:
        _, slide = new_slide()

        contained = core.add_image(
            slide,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(4),
            SAMPLE_IMAGE,
        )
        self.assertEqual(contained.width, Inches(4))
        self.assertEqual(contained.height, Inches(2))
        self.assertEqual(contained.left, Inches(1))
        self.assertEqual(contained.top, Inches(2))

        stretched = core.add_image(
            slide,
            Inches(6),
            Inches(1),
            Inches(3),
            Inches(2),
            SAMPLE_IMAGE,
            fit="stretch",
        )
        self.assertEqual(stretched.width, Inches(3))
        self.assertEqual(stretched.height, Inches(2))
        self.assertEqual(stretched.left, Inches(6))
        self.assertEqual(stretched.top, Inches(1))

        with self.assertRaisesRegex(FileNotFoundError, "missing.png"):
            core.add_image(slide, 0, 0, 1, 1, "missing.png")

        with self.assertRaisesRegex(
            ValueError, "image fit must be 'contain' or 'stretch'"
        ):
            core.add_image(slide, 0, 0, 1, 1, SAMPLE_IMAGE, fit="cover")

    def test_add_table_variants_and_validation(self) -> None:
        _, slide = new_slide()

        frame = core.add_table(
            slide,
            Inches(0.5),
            Inches(0.5),
            Inches(6),
            Inches(2),
            rows=[["_id", "ObjectId"], ["createdAt", "datetime"]],
            columns=["Field", "Type"],
            column_widths=[2.5, 3.5],
            row_heights=[0.5, 0.6, 0.6],
            banded_rows=True,
            style={"padding": "4pt", "banded-row-fill-color": "#EEEEEE"},
            header_style={"fill-color": "#193952", "font-color": "#FFFFFF"},
            cell_style={"font-color": "#111111"},
        )

        table = frame.table
        self.assertTrue(table.first_row)
        self.assertTrue(table.horz_banding)
        self.assertEqual(table.cell(0, 0).text, "Field")
        self.assertEqual(table.cell(1, 0).text, "_id")
        self.assertEqual(table.columns[0].width, Inches(2.5))
        self.assertEqual(table.rows[0].height, Inches(0.5))
        self.assertEqual(rgb_hex(table.cell(0, 0).fill.fore_color.rgb), "193952")
        self.assertEqual(rgb_hex(table.cell(2, 0).fill.fore_color.rgb), "EEEEEE")

        with self.assertRaisesRegex(TypeError, "table rows must be lists"):
            core.add_table(
                slide,
                0,
                0,
                1,
                1,
                rows=["bad"],  # type: ignore[list-item]
                columns=["A"],
            )

        with self.assertRaisesRegex(
            ValueError, "table must define at least one column"
        ):
            core.add_table(slide, 0, 0, 1, 1, rows=[], columns=None)

        with self.assertRaisesRegex(
            ValueError, "table rows must all match the number of columns"
        ):
            core.add_table(
                slide,
                0,
                0,
                1,
                1,
                rows=[["a", "b"]],
                columns=["A"],
            )

        with self.assertRaisesRegex(
            ValueError, "column_widths must match the number of columns"
        ):
            core.add_table(
                slide,
                0,
                0,
                1,
                1,
                rows=[["a"]],
                columns=["A"],
                column_widths=[1, 2],
            )

        with self.assertRaisesRegex(
            ValueError, "column_widths entries must be positive lengths"
        ):
            core.add_table(
                slide,
                0,
                0,
                1,
                1,
                rows=[["a"]],
                columns=["A"],
                column_widths=[0],
            )

        with self.assertRaisesRegex(
            ValueError,
            "row_heights must match the total number of header and body rows",
        ):
            core.add_table(
                slide,
                0,
                0,
                1,
                1,
                rows=[["a"]],
                columns=["A"],
                row_heights=[1],
            )

        with self.assertRaisesRegex(
            ValueError, "row_heights entries must be positive lengths"
        ):
            core.add_table(
                slide,
                0,
                0,
                1,
                1,
                rows=[["a"]],
                columns=["A"],
                row_heights=[1, 0],
            )

    def test_add_flow_boxes_set_notes_and_placeholder_helpers(self) -> None:
        _, slide = new_slide()
        core.add_flow_boxes(slide, [])
        self.assertEqual(len(slide.shapes), 0)

        core.add_flow_boxes(
            slide,
            [
                {
                    "label": "Extract",
                    "desc": "Get docs",
                    "style": {"fill-color": "#2E86AB", "arrow-color": "#654321"},
                },
                {
                    "label": "Load",
                    "desc": "Write docs",
                    "color": "#E86F51",
                },
            ],
        )
        texts = "\n".join(slide_texts(slide))
        self.assertIn("Extract", texts)
        self.assertIn("Get docs", texts)
        self.assertIn("→", texts)

        core.set_notes(slide, "Speaker note")
        self.assertEqual(slide.notes_slide.notes_text_frame.text, "Speaker note")

    def test_replace_placeholders_and_remove_generated_placeholders(self) -> None:
        _, slide = new_slide()
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(4), Inches(0.8)
        )
        title.text = "{{Title}}"
        bullets = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(4), Inches(1.5)
        )
        bullets.text = "{{ITEMS}}"
        cleared = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(4), Inches(0.8)
        )
        cleared.text = "{{empty}}"

        core.replace_placeholders(
            slide,
            {
                "title": "MongoDB",
                "items": ["ObjectId", "datetime"],
                "empty": None,
            },
            styles={
                ".slide": {"font-size": 24},
                "#title": {"uppercase": True},
            },
        )

        texts = "\n".join(slide_texts(slide))
        self.assertIn("MONGODB", texts)
        self.assertIn("ObjectId", texts)
        self.assertIn("datetime", texts)
        self.assertEqual(cleared.text, "")

        placeholder = slide.shapes.add_textbox(
            core.CONTENT_LEFT,
            core.CONTENT_TOP,
            Inches(4),
            Inches(1),
        )
        placeholder.text = "Body Text"
        unresolved = slide.shapes.add_textbox(
            core.CONTENT_LEFT,
            core.CONTENT_TOP + Inches(1),
            Inches(4),
            Inches(1),
        )
        unresolved.text = "{{body}}"
        above_threshold = slide.shapes.add_textbox(
            core.CONTENT_LEFT,
            Inches(1),
            Inches(4),
            Inches(1),
        )
        above_threshold.text = "{{keep}}"

        core.remove_generated_content_placeholders(slide)
        remaining_texts = "\n".join(slide_texts(slide))
        self.assertNotIn("Body Text", remaining_texts)
        self.assertNotIn("{{body}}", remaining_texts)
        self.assertIn("{{keep}}", remaining_texts)

    def test_clone_delete_and_move_slide(self) -> None:
        prs = Presentation(str(TEMPLATE))
        original_count = len(prs.slides)
        template_slide = prs.slides[3]

        cloned = core.clone_slide(prs, 3)
        self.assertEqual(len(prs.slides), original_count + 1)
        self.assertEqual(len(clone_texts(template_slide)), len(clone_texts(cloned)))

        clone_index = len(prs.slides) - 1
        cloned_title = clone_texts(cloned)[0]
        core.move_slide(prs, clone_index, 0)
        self.assertEqual(clone_texts(prs.slides[0])[0], cloned_title)

        core.delete_slide(prs, 0)
        self.assertEqual(len(prs.slides), original_count)

    def test_layout_content_shapes_with_image_adds_picture(self) -> None:
        prs = Presentation(str(TEMPLATE))
        slide = prs.slides[3]
        core.layout_content_shapes(
            slide,
            image={"path": SAMPLE_IMAGE, "caption": "Sample image"},
        )

        self.assertTrue(
            any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
        )
        self.assertIn("Sample image", "\n".join(slide_texts(slide)))


def clone_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            texts.append(shape.text_frame.text)
    return texts


if __name__ == "__main__":
    unittest.main()
