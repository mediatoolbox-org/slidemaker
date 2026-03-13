"""
title: Core utilities for manipulating python-pptx presentations.
summary: |-
  Provides low-level helper functions used by the template module
  to find shapes, set text, add bullet lists, and clone slides.
"""

from __future__ import annotations

import copy
import re
from pathlib import Path
from typing import Any, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.parts.image import Image as PptxImage
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu


# ── Style constants ─────────────────────────────────────────────
FONT_NAME = "Montserrat"
FONT_COLOR = RGBColor(0x0B, 0x1F, 0x33)
TITLE_FONT_SIZE = Pt(51)
SUBTITLE_FONT_SIZE = Pt(36)
BODY_FONT_SIZE = Pt(30)
SLIDE_W = Emu(18288000)  # 20 inches
SLIDE_H = Emu(10287000)  # 11.25 inches

# Content area defaults (below the title + decorative line)
CONTENT_LEFT = Inches(0.9)
CONTENT_TOP = Inches(2.6)
CONTENT_WIDTH = Inches(12.5)
CONTENT_HEIGHT = Inches(7.5)

# Brand palette — reusable colours for flow boxes etc.
BRAND = RGBColor(0x19, 0x39, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
ACCENT2 = RGBColor(0xE8, 0x6F, 0x51)
GREEN = RGBColor(0x48, 0xA9, 0x9A)
CODE_FONT = "Consolas"
CODE_FONT_SIZE = Pt(20)

_ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}
_CODE_LINE_PREFIX_RE = re.compile(r"^\s{0,2}\d+\s{1,2}")
_INLINE_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_INLINE_MARKDOWN_RE = re.compile(r"(\*\*[^*\n]+\*\*|`[^`\n]+`|\*[^*\n]+\*)")
_MARKDOWN_HEADING_RE = re.compile(r"^(#{1,3})\s+(.*)$")
_MARKDOWN_BULLET_RE = re.compile(r"^(\s*)[-*]\s+(.*)$")


def _normalize_style(style: Optional[dict[str, Any]]) -> dict[str, Any]:
    """
    title: Normalize style keys to kebab-case lowercase.
    parameters:
      style:
        type: Optional[dict[str, Any]]
    returns:
      type: dict[str, Any]
    """
    if not style:
        return {}
    normalized: dict[str, Any] = {}
    for key, value in style.items():
        if not isinstance(key, str):
            continue
        normalized[key.strip().lower().replace("_", "-")] = value
    return normalized


def _merge_style(
    base: Optional[dict[str, Any]],
    override: Optional[dict[str, Any]],
) -> dict[str, Any]:
    """
    title: Return a shallow merge where ``override`` wins.
    parameters:
      base:
        type: Optional[dict[str, Any]]
      override:
        type: Optional[dict[str, Any]]
    returns:
      type: dict[str, Any]
    """
    merged = dict(base or {})
    if override:
        merged.update(_normalize_style(override))
    return merged


def _as_rgb_color(
    value: Any,
    default: Optional[RGBColor] = None,
) -> Optional[RGBColor]:
    """
    title: Parse an RGBColor from common style representations.
    parameters:
      value:
        type: Any
      default:
        type: Optional[RGBColor]
    returns:
      type: Optional[RGBColor]
    """
    if value is None:
        return default
    if isinstance(value, RGBColor):
        return value
    if isinstance(value, str):
        text = value.strip()
        if text.startswith("#"):
            text = text[1:]
        if len(text) == 6:
            try:
                return RGBColor(
                    int(text[0:2], 16),
                    int(text[2:4], 16),
                    int(text[4:6], 16),
                )
            except ValueError:
                return default
    if isinstance(value, (tuple, list)) and len(value) == 3:
        try:
            r, g, b = (int(v) for v in value)
        except (TypeError, ValueError):
            return default
        if all(0 <= v <= 255 for v in (r, g, b)):
            return RGBColor(r, g, b)
    return default


def _as_pt(value: Any, default: Optional[int] = None) -> Optional[int]:
    """
    title: Parse a point-sized value as EMU using ``Pt``.
    parameters:
      value:
        type: Any
      default:
        type: Optional[int]
    returns:
      type: Optional[int]
    """
    if value is None:
        return default
    if isinstance(value, (int, float)):
        return Pt(float(value))
    if isinstance(value, str):
        text = value.strip().lower()
        if text.endswith("pt"):
            text = text[:-2].strip()
        try:
            return Pt(float(text))
        except ValueError:
            return default
    return default


def _as_length(value: Any, default: Optional[int] = None) -> Optional[int]:
    """
    title: Parse a PowerPoint length into EMU.
    summary: |-
      Numeric values are treated as inches for layout-oriented inputs
      such as table column widths and row heights. Strings may use
      ``in``, ``pt``, or ``emu`` suffixes.
    parameters:
      value:
        type: Any
      default:
        type: Optional[int]
    returns:
      type: Optional[int]
    """
    if value is None:
        return default

    emu = getattr(value, "emu", None)
    if emu is not None:
        try:
            return int(emu)
        except (TypeError, ValueError):
            return default

    if isinstance(value, bool):
        return default
    if isinstance(value, int):
        return value
    if isinstance(value, float):
        return Inches(float(value))
    if isinstance(value, str):
        text = value.strip().lower()
        if text.endswith("pt"):
            text = text[:-2].strip()
            try:
                return Pt(float(text))
            except ValueError:
                return default
        if text.endswith("in"):
            text = text[:-2].strip()
            try:
                return Inches(float(text))
            except ValueError:
                return default
        if text.endswith("emu"):
            text = text[:-3].strip()
            try:
                return int(text)
            except ValueError:
                return default
        try:
            return Inches(float(text))
        except ValueError:
            return default
    return default


def _font_size_pt(font_size: Any) -> Optional[float]:
    """
    title: Extract font size in points from a ``Length`` value.
    parameters:
      font_size:
        type: Any
    returns:
      type: Optional[float]
    """
    if font_size is None:
        return None
    pt = getattr(font_size, "pt", None)
    if pt is not None:
        try:
            return float(pt)
        except (TypeError, ValueError):
            return None
    return None


def _resolve_line_spacing(
    value: Any,
    font_size: Any,
    default: Any = None,
) -> Any:
    """
    title: Resolve line spacing for Canva-friendly output.
    summary: |-
      Numeric values are treated as multipliers and converted
      to fixed point leading based on font size. Use ``"pt"``
      suffix for absolute point values.
    parameters:
      value:
        type: Any
      font_size:
        type: Any
      default:
        type: Any
    returns:
      type: Any
    """
    if value is None:
        return default

    def as_points(multiplier: float) -> Any:
        font_pt = _font_size_pt(font_size)
        if font_pt is None:
            return multiplier
        return Pt(font_pt * multiplier)

    if isinstance(value, (int, float)):
        return as_points(float(value))

    if isinstance(value, str):
        text = value.strip().lower()
        if text.endswith("pt"):
            text = text[:-2].strip()
            try:
                return Pt(float(text))
            except ValueError:
                return default
        if text.endswith("x"):
            text = text[:-1].strip()
            try:
                return as_points(float(text))
            except ValueError:
                return default
        if text.endswith("%"):
            text = text[:-1].strip()
            try:
                return as_points(float(text) / 100.0)
            except ValueError:
                return default
        try:
            return as_points(float(text))
        except ValueError:
            return default

    return default


def _resolve_letter_spacing(
    value: Any,
    font_size: Any,
    default: Optional[int] = None,
) -> Optional[int]:
    """
    title: Resolve letter spacing for ``a:rPr@spc``.
    summary: |-
      Numeric values are treated as tracking units relative to
      font size (Canva style). ``"pt"`` values are absolute.
    parameters:
      value:
        type: Any
      font_size:
        type: Any
      default:
        type: Optional[int]
    returns:
      type: Optional[int]
    """
    if value is None:
        return default

    def tracking_to_spc(tracking: float) -> Optional[int]:
        font_pt = _font_size_pt(font_size)
        if font_pt is None:
            return int(round(tracking))
        # Canva-style tracking (per-thousand of em) to OOXML ST_TextPoint.
        return int(round(tracking * font_pt / 10.0))

    if isinstance(value, (int, float)):
        return tracking_to_spc(float(value))

    if isinstance(value, str):
        text = value.strip().lower()
        if text.endswith("pt"):
            text = text[:-2].strip()
            try:
                return int(round(float(text) * 100))
            except ValueError:
                return default
        try:
            return tracking_to_spc(float(text))
        except ValueError:
            return default

    return default


def _as_bool(value: Any, default: Optional[bool] = None) -> Optional[bool]:
    """
    title: Parse a boolean from common string/int representations.
    parameters:
      value:
        type: Any
      default:
        type: Optional[bool]
    returns:
      type: Optional[bool]
    """
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, int):
        return bool(value)
    if isinstance(value, str):
        text = value.strip().lower()
        if text in {"true", "1", "yes", "on"}:
            return True
        if text in {"false", "0", "no", "off"}:
            return False
    return default


def _resolve_uppercase(
    normalized: dict[str, Any],
    default: bool = False,
) -> bool:
    """
    title: Resolve uppercase transform from style keys.
    parameters:
      normalized:
        type: dict[str, Any]
      default:
        type: bool
    returns:
      type: bool
    """
    text_transform = normalized.get("text-transform")
    if isinstance(text_transform, str):
        mode = text_transform.strip().lower()
        if mode == "uppercase":
            return True
        if mode in {"none", "normal", "initial"}:
            return False
    return bool(_as_bool(normalized.get("uppercase"), default))


def _apply_uppercase(text: str, uppercase: bool) -> str:
    """
    title: Apply uppercase transform when enabled.
    parameters:
      text:
        type: str
      uppercase:
        type: bool
    returns:
      type: str
    """
    return text.upper() if uppercase else text


def _as_alignment(value: Any, default: Optional[int] = None) -> Optional[int]:
    """
    title: Parse paragraph alignment from text/int values.
    parameters:
      value:
        type: Any
      default:
        type: Optional[int]
    returns:
      type: Optional[int]
    """
    if value is None:
        return default
    if isinstance(value, int):
        return value
    if isinstance(value, str):
        return _ALIGNMENTS.get(value.strip().lower(), default)
    return default


def _with_code_line_numbers(code_text: str) -> str:
    """
    title: Prefix code lines with right-aligned line numbers.
    summary: |-
      Uses ``" X  "`` for one-digit lines and ``"XX  "`` for
      two-digit lines. If all non-empty lines already look
      numbered, the input is returned unchanged.
    parameters:
      code_text:
        type: str
    returns:
      type: str
    """
    lines = code_text.splitlines()
    non_empty = [line for line in lines if line.strip()]
    if non_empty and all(_CODE_LINE_PREFIX_RE.match(line) for line in non_empty):
        return code_text
    return "\n".join(f"{idx:>2}  {line}" for idx, line in enumerate(lines, 1))


def _markdown_bold_segments(text: str) -> list[tuple[str, bool]]:
    """
    title: Split text into plain/bold segments based on ``**...**`` markup.
    parameters:
      text:
        type: str
    returns:
      type: list[tuple[str, bool]]
    """
    segments: list[tuple[str, bool]] = []
    last = 0
    for match in _INLINE_BOLD_RE.finditer(text):
        if match.start() > last:
            plain = text[last : match.start()].replace("**", "")
            if plain:
                segments.append((plain, False))
        bold_text = match.group(1)
        if bold_text:
            segments.append((bold_text, True))
        last = match.end()

    tail = text[last:].replace("**", "")
    if tail:
        segments.append((tail, False))

    if not segments:
        clean = text.replace("**", "")
        if clean:
            segments.append((clean, False))
    return segments


def _markdown_inline_segments(text: str) -> list[tuple[str, bool, bool, bool]]:
    """
    title: Split inline markdown into styled text segments.
    summary: |-
      Supports a small inline subset used in generated text blocks:
      ``**bold**``, ``*italic*``, and `` `code` ``.
    parameters:
      text:
        type: str
    returns:
      type: list[tuple[str, bool, bool, bool]]
    """
    segments: list[tuple[str, bool, bool, bool]] = []
    last = 0
    for match in _INLINE_MARKDOWN_RE.finditer(text):
        if match.start() > last:
            plain = text[last : match.start()]
            if plain:
                segments.append((plain, False, False, False))

        token = match.group(0)
        if token.startswith("**") and token.endswith("**"):
            segments.append((token[2:-2], True, False, False))
        elif token.startswith("`") and token.endswith("`"):
            segments.append((token[1:-1], False, False, True))
        else:
            segments.append((token[1:-1], False, True, False))
        last = match.end()

    tail = text[last:]
    if tail:
        segments.append((tail, False, False, False))

    if not segments and text:
        segments.append((text, False, False, False))
    return segments


def _resolve_padding(
    normalized: dict[str, Any],
    default: Optional[int] = None,
) -> tuple[Optional[int], Optional[int], Optional[int], Optional[int]]:
    """
    title: Resolve padding values (left, top, right, bottom) in EMU.
    parameters:
      normalized:
        type: dict[str, Any]
      default:
        type: Optional[int]
    returns:
      type: tuple[Optional[int], Optional[int], Optional[int], Optional[int]]
    """
    pad_all = _as_pt(normalized.get("padding"), default)
    pad_x = _as_pt(normalized.get("padding-x"), pad_all)
    pad_y = _as_pt(normalized.get("padding-y"), pad_all)

    pad_left = _as_pt(normalized.get("padding-left"), pad_x)
    pad_top = _as_pt(normalized.get("padding-top"), pad_y)
    pad_right = _as_pt(normalized.get("padding-right"), pad_x)
    pad_bottom = _as_pt(normalized.get("padding-bottom"), pad_y)
    return pad_left, pad_top, pad_right, pad_bottom


def _apply_text_frame_padding(
    tf: Any,
    normalized: dict[str, Any],
    default: Optional[int] = None,
) -> None:
    """
    title: Apply resolved padding values to a text frame margins.
    parameters:
      tf:
        type: Any
      normalized:
        type: dict[str, Any]
      default:
        type: Optional[int]
    """
    pad_left, pad_top, pad_right, pad_bottom = _resolve_padding(normalized, default)
    if pad_left is not None:
        tf.margin_left = pad_left
    if pad_top is not None:
        tf.margin_top = pad_top
    if pad_right is not None:
        tf.margin_right = pad_right
    if pad_bottom is not None:
        tf.margin_bottom = pad_bottom


def _apply_run_letter_spacing(run: Any, spacing: Optional[int]) -> None:
    """
    title: Apply letter spacing to a run via ``a:rPr@spc`` (centipoints).
    parameters:
      run:
        type: Any
      spacing:
        type: Optional[int]
    """
    if spacing is None:
        return
    r_pr = run._r.get_or_add_rPr()
    r_pr.set("spc", str(spacing))


def _apply_paragraph_bullet(
    paragraph: Any,
    bullet_char: str,
    level: int = 0,
) -> None:
    """
    title: Apply bullet formatting to a paragraph with hanging indentation.
    parameters:
      paragraph:
        type: Any
      bullet_char:
        type: str
      level:
        type: int
    """
    normalized_level = max(0, min(int(level), 8))
    left_margin = int(Pt(24 + (normalized_level * 18)))
    hanging_indent = -int(Pt(18))

    pPr = paragraph._p.get_or_add_pPr()
    for child in list(pPr):
        if child.tag.endswith(("}buNone", "}buAutoNum", "}buChar", "}buBlip")):
            pPr.remove(child)
    pPr.set("lvl", str(normalized_level))
    pPr.set("marL", str(left_margin))
    pPr.set("indent", str(hanging_indent))

    bu_char = OxmlElement("a:buChar")
    bu_char.set("char", bullet_char)
    pPr.append(bu_char)


def find_group_textbox(slide: Slide, group_name: str) -> Any:
    """
    title: Find the first TextBox inside a named Group shape.
    parameters:
      slide:
        type: Slide
        description: The slide to search.
      group_name:
        type: str
        description: The ``name`` attribute of the Group shape.
    returns:
      type: Any
      description: The TextBox shape if found, otherwise ``None``.
    """
    for shape in slide.shapes:
        if shape.name == group_name and shape.shape_type == 6:
            for child in shape.shapes:  # type: ignore[attr-defined]
                if (
                    child.has_text_frame and child.shape_type == 17  # TEXT_BOX
                ):
                    return child
    return None


def find_textbox_by_name(slide: Slide, name: str) -> Any:
    """
    title: Find a shape by its exact name on a slide.
    parameters:
      slide:
        type: Slide
        description: The slide to search.
      name:
        type: str
        description: The ``name`` attribute of the shape.
    returns:
      type: Any
      description: The shape if found, otherwise ``None``.
    """
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_textbox_text(
    shape: Any,
    text: str,
    font_size: Optional[int] = None,
    font_color: Optional[RGBColor] = None,
    font_name: Optional[str] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    alignment: Optional[int] = None,
    style: Optional[dict[str, Any]] = None,
) -> None:
    """
    title: Replace all text in a shape's text frame.
    parameters:
      shape:
        type: Any
        description: A shape with a ``text_frame`` attribute.
      text:
        type: str
        description: The replacement text.
      font_size:
        type: Optional[int]
        description: Font size in EMU (use ``Pt()``).
      font_color:
        type: Optional[RGBColor]
        description: Font colour.
      font_name:
        type: Optional[str]
        description: Font family name.
      bold:
        type: Optional[bool]
        description: Whether the text is bold.
      italic:
        type: Optional[bool]
      alignment:
        type: Optional[int]
        description: Paragraph alignment constant.
      style:
        type: Optional[dict[str, Any]]
    """
    normalized = _normalize_style(style)
    resolved_uppercase = _resolve_uppercase(normalized, False)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = _apply_uppercase(text, resolved_uppercase)
    _apply_text_frame_padding(tf, normalized)

    resolved_alignment = alignment
    if resolved_alignment is None:
        resolved_alignment = _as_alignment(
            normalized.get("alignment", normalized.get("align")),
        )
    if resolved_alignment is not None:
        p.alignment = resolved_alignment

    # python-pptx may produce zero runs when text is empty.
    if not p.runs:
        return

    resolved_font_size = (
        font_size if font_size is not None else _as_pt(normalized.get("font-size"))
    )
    resolved_font_color = (
        font_color
        if font_color is not None
        else _as_rgb_color(normalized.get("font-color"))
    )
    resolved_font_name = (
        font_name if font_name is not None else normalized.get("font-name")
    )
    if resolved_font_name is not None:
        resolved_font_name = str(resolved_font_name)
    resolved_bold = bold if bold is not None else _as_bool(normalized.get("bold"))
    resolved_italic = (
        italic if italic is not None else _as_bool(normalized.get("italic"))
    )
    resolved_line_spacing = _resolve_line_spacing(
        normalized.get("line-spacing", normalized.get("line-height")),
        resolved_font_size,
    )
    if resolved_line_spacing is not None:
        p.line_spacing = resolved_line_spacing
    resolved_letter_spacing = _resolve_letter_spacing(
        normalized.get("letter-spacing"),
        resolved_font_size,
    )

    run = p.runs[0]
    if resolved_font_size is not None:
        run.font.size = resolved_font_size
    if resolved_font_color is not None:
        run.font.color.rgb = resolved_font_color
    if resolved_font_name is not None:
        run.font.name = resolved_font_name
    if resolved_bold is not None:
        run.font.bold = resolved_bold
    if resolved_italic is not None:
        run.font.italic = resolved_italic
    _apply_run_letter_spacing(run, resolved_letter_spacing)


def add_textbox(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: Optional[int] = None,
    font_color: Optional[RGBColor] = None,
    font_name: Optional[str] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    alignment: Optional[int] = None,
    style: Optional[dict[str, Any]] = None,
) -> object:
    """
    title: Add a new text box to a slide.
    parameters:
      slide:
        type: Slide
        description: Target slide.
      left:
        type: int
        description: Horizontal position in EMU.
      top:
        type: int
        description: Vertical position in EMU.
      width:
        type: int
        description: Box width in EMU.
      height:
        type: int
        description: Box height in EMU.
      text:
        type: str
        description: The text content.
      font_size:
        type: Optional[int]
        description: Font size in EMU (use ``Pt()``).
      font_color:
        type: Optional[RGBColor]
        description: Font colour.
      font_name:
        type: Optional[str]
        description: Font family name.
      bold:
        type: Optional[bool]
        description: Whether the text is bold.
      italic:
        type: Optional[bool]
      alignment:
        type: Optional[int]
        description: Paragraph alignment constant.
      style:
        type: Optional[dict[str, Any]]
    returns:
      type: object
      description: The newly created text box shape.
    """
    normalized = _normalize_style(style)
    resolved_font_size = (
        font_size
        if font_size is not None
        else _as_pt(normalized.get("font-size"), BODY_FONT_SIZE)
    )
    resolved_font_color = (
        font_color
        if font_color is not None
        else _as_rgb_color(normalized.get("font-color"), FONT_COLOR)
    )
    resolved_font_name_raw = (
        font_name if font_name is not None else normalized.get("font-name", FONT_NAME)
    )
    resolved_font_name = (
        FONT_NAME if resolved_font_name_raw is None else str(resolved_font_name_raw)
    )
    resolved_bold = (
        bold if bold is not None else _as_bool(normalized.get("bold"), False)
    )
    resolved_italic = (
        italic if italic is not None else _as_bool(normalized.get("italic"), False)
    )
    resolved_alignment = (
        alignment
        if alignment is not None
        else _as_alignment(
            normalized.get("alignment", normalized.get("align")),
            PP_ALIGN.LEFT,
        )
    )
    resolved_line_spacing = _resolve_line_spacing(
        normalized.get("line-spacing", normalized.get("line-height")),
        resolved_font_size,
    )
    resolved_letter_spacing = _resolve_letter_spacing(
        normalized.get("letter-spacing"),
        resolved_font_size,
    )
    resolved_uppercase = _resolve_uppercase(normalized, False)

    txbox = slide.shapes.add_textbox(left, top, width, height)  # type: ignore[arg-type]
    tf = txbox.text_frame
    tf.word_wrap = True
    _apply_text_frame_padding(tf, normalized)
    p = tf.paragraphs[0]
    p.text = _apply_uppercase(text, resolved_uppercase)
    p.font.size = resolved_font_size
    p.font.color.rgb = resolved_font_color
    p.font.name = resolved_font_name
    p.font.bold = resolved_bold
    p.font.italic = resolved_italic
    p.alignment = resolved_alignment
    if resolved_line_spacing is not None:
        p.line_spacing = resolved_line_spacing
    for run in p.runs:
        _apply_run_letter_spacing(run, resolved_letter_spacing)
    return txbox


def add_markdown_textbox(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    markdown_text: str,
    style: Optional[dict[str, Any]] = None,
) -> object:
    """
    title: Add a free-form markdown text block to a slide.
    summary: |-
      Supports paragraphs, ``#``/``##``/``###`` headings, unordered
      list items using ``-`` or ``*`` as real PowerPoint bullets,
      nested bullets via two-space indentation, and inline
      ``**bold**``, ``*italic*``, and `` `code` `` markup.
    parameters:
      slide:
        type: Slide
      left:
        type: int
      top:
        type: int
      width:
        type: int
      height:
        type: int
      markdown_text:
        type: str
      style:
        type: Optional[dict[str, Any]]
    returns:
      type: object
      description: The newly created text box shape.
    """
    normalized = _normalize_style(style)
    resolved_font_size = _as_pt(normalized.get("font-size"), BODY_FONT_SIZE)
    resolved_font_color = _as_rgb_color(normalized.get("font-color"), FONT_COLOR)
    resolved_font_name_raw = normalized.get("font-name", FONT_NAME)
    resolved_font_name = (
        FONT_NAME if resolved_font_name_raw is None else str(resolved_font_name_raw)
    )
    resolved_bold = _as_bool(normalized.get("bold"), False)
    resolved_italic = _as_bool(normalized.get("italic"), False)
    resolved_alignment = _as_alignment(
        normalized.get("alignment", normalized.get("align")),
        PP_ALIGN.LEFT,
    )
    resolved_spacing = _as_pt(
        normalized.get("spacing", normalized.get("space-after")),
        Pt(10),
    )
    resolved_space_before = _as_pt(normalized.get("space-before"), Pt(0))
    resolved_line_spacing = _resolve_line_spacing(
        normalized.get("line-spacing", normalized.get("line-height")),
        resolved_font_size,
    )
    resolved_letter_spacing = _resolve_letter_spacing(
        normalized.get("letter-spacing"),
        resolved_font_size,
    )
    resolved_uppercase = _resolve_uppercase(normalized, False)
    resolved_bullet_char = str(normalized.get("bullet-char", "•"))

    txbox = slide.shapes.add_textbox(left, top, width, height)  # type: ignore[arg-type]
    tf = txbox.text_frame
    tf.word_wrap = True
    _apply_text_frame_padding(tf, normalized)

    lines = markdown_text.strip("\n").split("\n")
    first_paragraph = True
    paragraph_gap = False
    base_font_size_pt = _font_size_pt(resolved_font_size) or BODY_FONT_SIZE.pt

    for raw_line in lines:
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            paragraph_gap = True
            continue

        p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
        first_paragraph = False
        p.alignment = resolved_alignment
        p.space_after = resolved_spacing
        p.space_before = resolved_spacing if paragraph_gap else resolved_space_before
        if resolved_line_spacing is not None:
            p.line_spacing = resolved_line_spacing

        heading_match = _MARKDOWN_HEADING_RE.match(stripped)
        bullet_match = _MARKDOWN_BULLET_RE.match(line)
        heading_level: int | None = None
        bullet_level = 0
        paragraph_text = stripped
        if heading_match is not None:
            heading_level = min(3, len(heading_match.group(1)))
            paragraph_text = heading_match.group(2).strip()
        elif bullet_match is not None:
            leading_spaces = len(bullet_match.group(1).expandtabs(2))
            bullet_level = leading_spaces // 2
            paragraph_text = bullet_match.group(2).strip()
            _apply_paragraph_bullet(p, resolved_bullet_char, bullet_level)

        paragraph_text = _apply_uppercase(paragraph_text, resolved_uppercase)
        paragraph_font_size = resolved_font_size
        paragraph_bold = resolved_bold
        if heading_level is not None:
            heading_size_map = {
                1: Pt(base_font_size_pt + 8),
                2: Pt(base_font_size_pt + 4),
                3: Pt(base_font_size_pt + 2),
            }
            paragraph_font_size = heading_size_map[heading_level]
            paragraph_bold = True

        for segment_text, is_bold, is_italic, is_code in _markdown_inline_segments(
            paragraph_text
        ):
            run = p.add_run()
            run.text = segment_text
            run.font.size = paragraph_font_size
            run.font.color.rgb = resolved_font_color
            run.font.name = CODE_FONT if is_code else resolved_font_name
            run.font.bold = bool(paragraph_bold or is_bold)
            run.font.italic = bool(resolved_italic or is_italic)
            _apply_run_letter_spacing(run, resolved_letter_spacing)

        paragraph_gap = False

    return txbox


def add_bullet_list(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    items: list[str],
    font_size: Optional[int] = None,
    font_color: Optional[RGBColor] = None,
    font_name: Optional[str] = None,
    spacing: Optional[int] = None,
    bullet_char: Optional[str] = None,
    bold_prefixes: Optional[bool] = None,
    style: Optional[dict[str, Any]] = None,
) -> object:
    """
    title: Add a bulleted list as a text box on a slide.
    summary: |-
      Supports basic inline markdown bold: ``**text**`` segments
      render in bold when ``bold_prefixes`` is ``True``.
    parameters:
      slide:
        type: Slide
        description: Target slide.
      left:
        type: int
        description: Horizontal position in EMU.
      top:
        type: int
        description: Vertical position in EMU.
      width:
        type: int
        description: Box width in EMU.
      height:
        type: int
        description: Box height in EMU.
      items:
        type: list[str]
        description: Bullet point strings.
      font_size:
        type: Optional[int]
        description: Font size in EMU.
      font_color:
        type: Optional[RGBColor]
        description: Font colour.
      font_name:
        type: Optional[str]
        description: Font family name.
      spacing:
        type: Optional[int]
        description: Space after each paragraph in EMU.
      bullet_char:
        type: Optional[str]
        description: Character used for real paragraph bullets.
      bold_prefixes:
        type: Optional[bool]
        description: >-
          If ``True``, parse and render ``**...**`` segments in bold within
          each item.
      style:
        type: Optional[dict[str, Any]]
    returns:
      type: object
      description: The newly created text box shape.
    """
    normalized = _normalize_style(style)
    resolved_font_size = (
        font_size
        if font_size is not None
        else _as_pt(normalized.get("font-size"), BODY_FONT_SIZE)
    )
    resolved_font_color = (
        font_color
        if font_color is not None
        else _as_rgb_color(normalized.get("font-color"), FONT_COLOR)
    )
    resolved_font_name_raw = (
        font_name if font_name is not None else normalized.get("font-name", FONT_NAME)
    )
    resolved_font_name = (
        FONT_NAME if resolved_font_name_raw is None else str(resolved_font_name_raw)
    )
    resolved_spacing = (
        spacing
        if spacing is not None
        else _as_pt(
            normalized.get("spacing", normalized.get("space-after")),
            Pt(10),
        )
    )
    resolved_bullet_char = (
        bullet_char
        if bullet_char is not None
        else str(normalized.get("bullet-char", "•"))
    )
    resolved_bold_prefixes = (
        bold_prefixes
        if bold_prefixes is not None
        else _as_bool(normalized.get("bold-prefixes"), True)
    )
    resolved_alignment = _as_alignment(
        normalized.get("alignment", normalized.get("align")),
        PP_ALIGN.LEFT,
    )
    resolved_space_before = _as_pt(normalized.get("space-before"), Pt(2))
    resolved_italic = _as_bool(normalized.get("italic"), False)
    resolved_line_spacing = _resolve_line_spacing(
        normalized.get("line-spacing", normalized.get("line-height")),
        resolved_font_size,
    )
    resolved_letter_spacing = _resolve_letter_spacing(
        normalized.get("letter-spacing"),
        resolved_font_size,
    )
    resolved_uppercase = _resolve_uppercase(normalized, False)

    txbox = slide.shapes.add_textbox(left, top, width, height)  # type: ignore[arg-type]
    tf = txbox.text_frame
    tf.word_wrap = True
    _apply_text_frame_padding(tf, normalized)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = resolved_spacing
        p.space_before = resolved_space_before
        p.alignment = resolved_alignment
        if resolved_line_spacing is not None:
            p.line_spacing = resolved_line_spacing

        _apply_paragraph_bullet(p, resolved_bullet_char)

        item_text = _apply_uppercase(item.strip(), resolved_uppercase)
        if resolved_bold_prefixes:
            segments = _markdown_bold_segments(item_text)
        else:
            segments = [(item_text.replace("**", ""), False)]

        for segment_text, is_bold in segments:
            run = p.add_run()
            run.text = segment_text
            run.font.size = resolved_font_size
            run.font.color.rgb = resolved_font_color
            run.font.name = resolved_font_name
            run.font.bold = is_bold
            run.font.italic = resolved_italic
            _apply_run_letter_spacing(run, resolved_letter_spacing)

    return txbox


def add_shape_rect(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_color: Optional[RGBColor] = None,
    line_color: Optional[RGBColor] = None,
    line_width: Optional[int] = None,
    style: Optional[dict[str, Any]] = None,
) -> object:
    """
    title: Add a rectangle shape to a slide.
    parameters:
      slide:
        type: Slide
        description: Target slide.
      left:
        type: int
        description: Horizontal position in EMU.
      top:
        type: int
        description: Vertical position in EMU.
      width:
        type: int
        description: Shape width in EMU.
      height:
        type: int
        description: Shape height in EMU.
      fill_color:
        type: Optional[RGBColor]
        description: Fill colour. If ``None`` the shape has no fill.
      line_color:
        type: Optional[RGBColor]
      line_width:
        type: Optional[int]
      style:
        type: Optional[dict[str, Any]]
    returns:
      type: object
      description: The newly created rectangle shape.
    """
    normalized = _normalize_style(style)
    resolved_fill = (
        fill_color
        if fill_color is not None
        else _as_rgb_color(normalized.get("fill-color"))
    )
    resolved_line_color = (
        line_color
        if line_color is not None
        else _as_rgb_color(normalized.get("line-color"))
    )
    resolved_line_width = (
        line_width if line_width is not None else _as_pt(normalized.get("line-width"))
    )

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)  # type: ignore[arg-type]
    if resolved_line_color is not None:
        shp.line.color.rgb = resolved_line_color
        if resolved_line_width is not None:
            shp.line.width = resolved_line_width
    else:
        shp.line.fill.background()
    if resolved_fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = resolved_fill
    else:
        shp.fill.background()
    return shp


def _set_fill_color(target: Any, color: Optional[RGBColor]) -> None:
    """
    title: Apply a solid fill or clear the background fill.
    parameters:
      target:
        type: Any
      color:
        type: Optional[RGBColor]
    """
    if color is None:
        target.fill.background()
        return
    target.fill.solid()
    target.fill.fore_color.rgb = color


def add_code_block(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    code_text: str,
    bg_color: Optional[RGBColor] = None,
    font_size: Optional[int] = None,
    font_color: Optional[RGBColor] = None,
    font_name: Optional[str] = None,
    style: Optional[dict[str, Any]] = None,
) -> None:
    """
    title: Add a code block with a dark background to a slide.
    summary: |-
      Renders monospace text on a solid-colour rectangle.
      Style options include ``bg-color``, ``font-color``,
      ``font-name``, ``font-size``, ``line-numbers``,
      ``letter-spacing``, ``line-spacing`` (or ``line-height``),
      ``uppercase`` (or ``text-transform: uppercase``), and
      ``padding`` (plus side-specific variants).
    parameters:
      slide:
        type: Slide
        description: Target slide.
      left:
        type: int
        description: Horizontal position in EMU.
      top:
        type: int
        description: Vertical position in EMU.
      width:
        type: int
        description: Block width in EMU.
      height:
        type: int
        description: Block height in EMU.
      code_text:
        type: str
        description: The source code to display.
      bg_color:
        type: Optional[RGBColor]
        description: Background rectangle colour.
      font_size:
        type: Optional[int]
        description: Font size for the code text.
      font_color:
        type: Optional[RGBColor]
      font_name:
        type: Optional[str]
      style:
        type: Optional[dict[str, Any]]
    """
    normalized = _normalize_style(style)
    resolved_bg_color = (
        bg_color
        if bg_color is not None
        else _as_rgb_color(
            normalized.get("bg-color", normalized.get("fill-color")),
            BRAND,
        )
    )
    resolved_font_size = (
        font_size
        if font_size is not None
        else _as_pt(normalized.get("font-size"), CODE_FONT_SIZE)
    )
    resolved_font_color = (
        font_color
        if font_color is not None
        else _as_rgb_color(normalized.get("font-color"), WHITE)
    )
    resolved_font_name_raw = (
        font_name if font_name is not None else normalized.get("font-name", CODE_FONT)
    )
    resolved_font_name = (
        CODE_FONT if resolved_font_name_raw is None else str(resolved_font_name_raw)
    )
    resolved_bold = _as_bool(normalized.get("bold"), False)
    resolved_line_spacing = _resolve_line_spacing(
        normalized.get("line-spacing", normalized.get("line-height")),
        resolved_font_size,
    )
    resolved_letter_spacing = _resolve_letter_spacing(
        normalized.get("letter-spacing"),
        resolved_font_size,
    )
    resolved_uppercase = _resolve_uppercase(normalized, False)
    resolved_line_numbers = _as_bool(normalized.get("line-numbers"), False)
    if resolved_line_numbers:
        code_text = _with_code_line_numbers(code_text)

    shp = add_shape_rect(slide, left, top, width, height, fill_color=resolved_bg_color)
    shp.shadow.inherit = False  # type: ignore[attr-defined]

    pad_left, pad_top, pad_right, pad_bottom = _resolve_padding(
        normalized, Inches(0.25)
    )
    inner_width = max(1, width - (pad_left or 0) - (pad_right or 0))
    inner_height = max(1, height - (pad_top or 0) - (pad_bottom or 0))
    txbox = slide.shapes.add_textbox(
        left + (pad_left or 0),  # type: ignore[arg-type]
        top + (pad_top or 0),  # type: ignore[arg-type]
        inner_width,  # type: ignore[arg-type]
        inner_height,  # type: ignore[arg-type]
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_text.strip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        if resolved_line_spacing is not None:
            p.line_spacing = resolved_line_spacing
        run = p.add_run()
        run.text = _apply_uppercase(line, resolved_uppercase)
        run.font.size = resolved_font_size
        run.font.color.rgb = resolved_font_color
        run.font.name = resolved_font_name
        run.font.bold = resolved_bold
        _apply_run_letter_spacing(run, resolved_letter_spacing)


def add_image(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    image_file: str | Path,
    fit: str = "contain",
) -> object:
    """
    title: Add an image to a slide, scaled within a target rectangle.
    summary: |-
      ``fit="contain"`` preserves aspect ratio and centers the image
      inside the target rectangle. ``fit="stretch"`` fills the target
      rectangle exactly.
    parameters:
      slide:
        type: Slide
      left:
        type: int
      top:
        type: int
      width:
        type: int
      height:
        type: int
      image_file:
        type: str | Path
      fit:
        type: str
    returns:
      type: object
      description: The created picture shape.
    """
    image_path = Path(image_file)
    if not image_path.is_file():
        raise FileNotFoundError(str(image_path))

    fit_mode = fit.strip().lower()
    if fit_mode not in {"contain", "stretch"}:
        raise ValueError("image fit must be 'contain' or 'stretch'")

    if fit_mode == "stretch":
        return slide.shapes.add_picture(
            str(image_path),
            Emu(left),
            Emu(top),
            Emu(width),
            Emu(height),
        )

    image = PptxImage.from_file(str(image_path))
    pixel_width, pixel_height = image.size
    if pixel_width <= 0 or pixel_height <= 0:
        raise ValueError("image dimensions must be positive")

    scale = min(width / pixel_width, height / pixel_height)
    picture_width = max(1, int(round(pixel_width * scale)))
    picture_height = max(1, int(round(pixel_height * scale)))
    picture_left = left + max(0, int(round((width - picture_width) / 2)))
    picture_top = top + max(0, int(round((height - picture_height) / 2)))

    return slide.shapes.add_picture(
        str(image_path),
        Emu(picture_left),
        Emu(picture_top),
        Emu(picture_width),
        Emu(picture_height),
    )


def add_table(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    rows: list[list[Any]],
    columns: Optional[list[Any]] = None,
    column_widths: Optional[list[Any]] = None,
    row_heights: Optional[list[Any]] = None,
    banded_rows: Optional[bool] = None,
    style: Optional[dict[str, Any]] = None,
    header_style: Optional[dict[str, Any]] = None,
    cell_style: Optional[dict[str, Any]] = None,
) -> object:
    """
    title: Add a table shape to a slide.
    summary: |-
      The table may include an optional header row via ``columns``.
      Table-level style keys apply to all cells. ``header_style`` and
      ``cell_style`` provide more specific text/fill overrides.
      Numeric ``column_widths`` and ``row_heights`` are interpreted as inches.
    parameters:
      slide:
        type: Slide
      left:
        type: int
      top:
        type: int
      width:
        type: int
      height:
        type: int
      rows:
        type: list[list[Any]]
      columns:
        type: Optional[list[Any]]
      column_widths:
        type: Optional[list[Any]]
      row_heights:
        type: Optional[list[Any]]
      banded_rows:
        type: Optional[bool]
      style:
        type: Optional[dict[str, Any]]
      header_style:
        type: Optional[dict[str, Any]]
      cell_style:
        type: Optional[dict[str, Any]]
    returns:
      type: object
      description: The created table graphic frame.
    """
    normalized = _normalize_style(style)
    header_values = ["" if value is None else str(value) for value in (columns or [])]

    body_rows: list[list[str]] = []
    for row in rows:
        if not isinstance(row, list):
            raise TypeError("table rows must be lists")
        body_rows.append(["" if value is None else str(value) for value in row])

    col_count = len(header_values)
    if col_count == 0 and body_rows:
        col_count = len(body_rows[0])
    if col_count == 0:
        raise ValueError("table must define at least one column")

    if header_values and len(header_values) != col_count:
        raise ValueError("table columns must all have the same length")

    for row in body_rows:
        if len(row) != col_count:
            raise ValueError("table rows must all match the number of columns")

    header_row_count = 1 if header_values else 0
    total_rows = len(body_rows) + header_row_count
    if total_rows == 0:
        raise ValueError("table must contain at least one row")

    frame = slide.shapes.add_table(
        total_rows,
        col_count,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    table = frame.table
    table.first_row = bool(header_values)
    resolved_banded_rows = (
        _as_bool(banded_rows, False)
        if banded_rows is not None
        else _as_bool(normalized.get("banded-rows"), False)
    )
    table.horz_banding = bool(resolved_banded_rows)

    if column_widths is not None:
        if len(column_widths) != col_count:
            raise ValueError("column_widths must match the number of columns")
        for idx, raw_width in enumerate(column_widths):
            resolved_width = _as_length(raw_width)
            if resolved_width is None or resolved_width <= 0:
                raise ValueError("column_widths entries must be positive lengths")
            table.columns[idx].width = resolved_width

    if row_heights is not None:
        if len(row_heights) != total_rows:
            raise ValueError(
                "row_heights must match the total number of header and body rows"
            )
        for idx, raw_height in enumerate(row_heights):
            resolved_height = _as_length(raw_height)
            if resolved_height is None or resolved_height <= 0:
                raise ValueError("row_heights entries must be positive lengths")
            table.rows[idx].height = resolved_height

    table_style = dict(normalized)
    effective_header_style = _merge_style(
        table_style,
        {
            "bold": True,
            "alignment": "center",
            "fill-color": "#193952",
            "font-color": "#FFFFFF",
        },
    )
    effective_header_style = _merge_style(effective_header_style, header_style)
    effective_cell_style = _merge_style(table_style, cell_style)

    default_body_fill = _as_rgb_color(effective_cell_style.get("fill-color"))
    banded_row_fill = _as_rgb_color(
        effective_cell_style.get("banded-row-fill-color")
        or table_style.get("banded-row-fill-color")
    )

    def apply_cell_text(cell: Any, value: str, style_map: dict[str, Any]) -> None:
        cell.text_frame.word_wrap = True
        set_textbox_text(cell, value, style=style_map)
        _set_fill_color(cell, _as_rgb_color(style_map.get("fill-color")))

    row_idx = 0
    if header_values:
        for col_idx, value in enumerate(header_values):
            apply_cell_text(table.cell(row_idx, col_idx), value, effective_header_style)
        row_idx += 1

    for body_idx, row_values in enumerate(body_rows):
        row_fill = default_body_fill
        if row_fill is None and resolved_banded_rows and banded_row_fill is not None:
            if body_idx % 2 == 1:
                row_fill = banded_row_fill
        for col_idx, value in enumerate(row_values):
            cell = table.cell(row_idx + body_idx, col_idx)
            cell.text_frame.word_wrap = True
            set_textbox_text(cell, value, style=effective_cell_style)
            _set_fill_color(cell, row_fill)

    return frame


def add_flow_boxes(
    slide: Slide,
    boxes: list[dict],
    left: int = CONTENT_LEFT,
    top: int = CONTENT_TOP,
    box_width: Optional[int] = None,
    box_height: int = Inches(2.4),
    gap: int = Inches(0.5),
    style: Optional[dict[str, Any]] = None,
) -> None:
    """
    title: Add a horizontal flow diagram with coloured boxes and arrows.
    summary: |-
      Each box is a dictionary with keys:

      - ``label`` (str): bold heading inside the box.
      - ``desc`` (str): smaller description text (may contain
      newlines).
      - ``style`` (dict): style attributes for this box
      (for example ``fill-color``/``font-color``).
      - ``color`` (str or RGBColor): legacy fill colour key.
    parameters:
      slide:
        type: Slide
        description: Target slide.
      boxes:
        type: list[dict]
        description: Box definitions (see above).
      left:
        type: int
        description: Horizontal start position in EMU.
      top:
        type: int
        description: Vertical position in EMU.
      box_width:
        type: Optional[int]
        description: >-
          Width of each box.  If ``None`` it is computed to fill the available
          content width.
      box_height:
        type: int
        description: Height of each box.
      gap:
        type: int
        description: Space between boxes (includes arrow).
      style:
        type: Optional[dict[str, Any]]
        description: Default style for all boxes in the flow.
    """
    n = len(boxes)
    if n == 0:
        return

    if box_width is None:
        avail = CONTENT_WIDTH - gap * (n - 1)
        box_width = int(avail / n)

    base_style = _normalize_style(style)

    for i, box in enumerate(boxes):
        x = left + i * (box_width + gap)
        box_style = _merge_style(base_style, box.get("style"))
        color = _as_rgb_color(
            box_style.get("fill-color", box_style.get("color", box.get("color"))),
            ACCENT,
        )
        label_style = _merge_style(
            {
                "font-size": 28,
                "font-color": "#FFFFFF",
                "bold": True,
                "alignment": "center",
            },
            box_style,
        )
        desc_style = _merge_style(
            {
                "font-size": 22,
                "font-color": "#FFFFFF",
                "bold": False,
                "alignment": "center",
            },
            box_style,
        )
        arrow_style = _merge_style(
            {
                "font-size": 44,
                "font-color": "#193952",
                "bold": True,
                "alignment": "center",
            },
            box_style,
        )
        if "arrow-color" in box_style:
            arrow_style["font-color"] = box_style["arrow-color"]
        if "arrow-font-size" in box_style:
            arrow_style["font-size"] = box_style["arrow-font-size"]

        shp = add_shape_rect(
            slide,
            x,
            top,
            box_width,
            box_height,
            fill_color=color,
        )
        shp.shadow.inherit = False  # type: ignore[attr-defined]

        # Label
        add_textbox(
            slide,
            x,
            top + Inches(0.3),
            box_width,
            Inches(0.7),
            box["label"],
            style=label_style,
        )
        # Description
        if box.get("desc"):
            add_textbox(
                slide,
                x + Inches(0.2),
                top + Inches(1.1),
                box_width - Inches(0.4),
                box_height - Inches(1.3),
                box["desc"],
                style=desc_style,
            )

        # Arrow between boxes
        if i < n - 1:
            ax = x + box_width + Inches(0.02)
            add_textbox(
                slide,
                ax,
                top + Inches(0.6),
                gap - Inches(0.04),
                Inches(0.7),
                "→",
                style=arrow_style,
            )


def set_notes(slide: Slide, text: str) -> None:
    """
    title: Set the speaker notes for a slide.
    parameters:
      slide:
        type: Slide
        description: Target slide.
      text:
        type: str
        description: The notes content.
    """
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    if tf is not None:
        tf.text = text


_PLACEHOLDER_RE = re.compile(r"\{\{(\w+)\}\}")
_GENERATED_CONTENT_PLACEHOLDERS = {
    "content body placeholder",
    "body text",
    "click to add text",
}


def _iter_text_shapes(slide: Slide) -> list[Any]:
    """
    title: Yield all shapes (including group children) that have a text frame.
    parameters:
      slide:
        type: Slide
    returns:
      type: list[Any]
    """
    result: list[Any] = []
    for shape in slide.shapes:
        if shape.shape_type == 6:  # GROUP
            for child in shape.shapes:  # type: ignore[attr-defined]
                if child.has_text_frame:
                    result.append(child)
        elif shape.has_text_frame:
            result.append(shape)
    return result


def replace_placeholders(
    slide: Slide,
    content: dict[str, str | list[str] | None],
    styles: Optional[dict[str, dict[str, Any]]] = None,
) -> None:
    """
    title: Scan slide shapes for ``{{key}}`` placeholders and replace them.
    summary: |-
      Iterates every text shape (including children of Group shapes).
      When a shape's text matches ``{{key}}`` and *key* exists in
      *content*, the shape content is replaced:

      - ``str`` value  → replace the text.
      - ``list[str]``  → replace with a bullet list in the same
        position and size as the original shape.
      - ``None``       → clear the shape text.

      Per-placeholder styling is resolved from *styles* using the
      ``#key`` convention.  The ``.slide`` entry is used as the
      fallback for any placeholder without a dedicated ``#key`` style.
    parameters:
      slide:
        type: Slide
      content:
        type: dict[str, str | list[str] | None]
      styles:
        type: Optional[dict[str, dict[str, Any]]]
        description: >-
          Style map with ``.slide`` as base style and ``#key`` entries for per-
          placeholder overrides.
    """
    # Build case-insensitive lookup: lowercase key → original value
    lower_content: dict[str, str | list[str] | None] = {
        k.lower(): v for k, v in content.items()
    }
    base_style = dict((styles or {}).get(".slide", {}))

    def _style_for(key: str) -> dict[str, Any]:
        # Try exact #key, then lowercase #key
        override = (styles or {}).get(f"#{key}") or (styles or {}).get(
            f"#{key.lower()}"
        )
        if override is None:
            return base_style
        merged = dict(base_style)
        merged.update(override)
        return merged

    for shape in _iter_text_shapes(slide):
        text = shape.text_frame.text.strip()
        m = _PLACEHOLDER_RE.fullmatch(text)
        if m is None:
            continue
        key = m.group(1)
        lower_key = key.lower()
        if lower_key not in lower_content:
            continue
        value = lower_content[lower_key]
        effective_style = _style_for(key)
        if value is None:
            set_textbox_text(shape, "")
        elif isinstance(value, list):
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            add_bullet_list(
                slide,
                left=left,
                top=top,
                width=width,
                height=height,
                items=value,
                style=effective_style,
            )
        else:
            set_textbox_text(shape, str(value), style=effective_style)


def remove_generated_content_placeholders(slide: Slide) -> None:
    """
    title: Remove stock body placeholders before adding generated content.
    summary: |-
      Generated content such as bullets, code blocks, and tables is
      intended to own the main content area. This helper strips
      unresolved template placeholder text from that area to avoid
      visual overlap.
    parameters:
      slide:
        type: Slide
    """
    content_threshold = CONTENT_TOP - Inches(0.3)

    for shape in list(_iter_text_shapes(slide)):
        text = shape.text_frame.text.strip()
        if not text:
            continue
        lower_text = text.lower()
        is_unresolved_placeholder = _PLACEHOLDER_RE.fullmatch(text) is not None
        is_stock_placeholder = lower_text in _GENERATED_CONTENT_PLACEHOLDERS
        if not is_unresolved_placeholder and not is_stock_placeholder:
            continue
        if getattr(shape, "top", 0) < content_threshold:
            continue
        sp = shape._element
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)


def _table_spec_value(spec: dict[str, Any], *keys: str, default: Any = None) -> Any:
    """
    title: Read a table-spec option supporting hyphen and underscore spellings.
    parameters:
      spec:
        type: dict[str, Any]
      default:
        type: Any
      keys:
        type: str
        variadic: positional
    returns:
      type: Any
    """
    for key in keys:
        if key in spec:
            return spec[key]
        alternate = key.replace("-", "_") if "-" in key else key.replace("_", "-")
        if alternate in spec:
            return spec[alternate]
    return default


def _split_content_height(
    total_height: int,
    ratio: float,
    gap: int,
    min_first: int,
    min_second: int,
) -> tuple[int, int]:
    """
    title: Split a vertical layout area into two non-overlapping sections.
    parameters:
      total_height:
        type: int
      ratio:
        type: float
      gap:
        type: int
      min_first:
        type: int
      min_second:
        type: int
    returns:
      type: tuple[int, int]
    """
    available = max(1, total_height - gap)
    desired_first = int(round(available * ratio))
    max_first = max(1, available - min_second)
    if max_first < min_first:
        first = max(1, min(available - 1, desired_first))
    else:
        first = min(max_first, max(min_first, desired_first))
    second = max(1, available - first)
    return first, second


def layout_content_shapes(
    slide: Slide,
    items: list[str] | None = None,
    markdown: str | None = None,
    code: str | None = None,
    table: Optional[dict[str, Any]] = None,
    image: str | Path | dict[str, Any] | None = None,
    flow_boxes: list[dict] | None = None,
    callout: str | None = None,
    slide_style: Optional[dict[str, Any]] = None,
    code_style: Optional[dict[str, Any]] = None,
    table_style: Optional[dict[str, Any]] = None,
    table_header_style: Optional[dict[str, Any]] = None,
    table_cell_style: Optional[dict[str, Any]] = None,
) -> None:
    """
    title: Add content shapes to a slide with smart layout.
    summary: |-
      Places shapes in the content area below the title region,
      stacking vertically based on what is provided:

      - **flow_boxes** — flow diagram at the top of content area.
      - **items + code** — bullets on top, code below.
      - **markdown + code** — markdown text on top, code below.
      - **items + table** — bullets on top, table below.
      - **markdown + table** — markdown text on top, table below.
      - **code + table** — code on top, table below.
      - **items + image** — bullets on top, image below.
      - **markdown + image** — markdown text on top, image below.
      - **code + image** — code on top, image below.
      - **items only** — full content area.
      - **markdown only** — full content area.
      - **code only** — full content area.
      - **table only** — full content area.
      - **image only** — full content area.
      - **callout** — always placed below other content.
    parameters:
      slide:
        type: Slide
      items:
        type: list[str] | None
      markdown:
        type: str | None
      code:
        type: str | None
      table:
        type: Optional[dict[str, Any]]
      image:
        type: str | Path | dict[str, Any] | None
      flow_boxes:
        type: list[dict] | None
      callout:
        type: str | None
      slide_style:
        type: Optional[dict[str, Any]]
      code_style:
        type: Optional[dict[str, Any]]
      table_style:
        type: Optional[dict[str, Any]]
      table_header_style:
        type: Optional[dict[str, Any]]
      table_cell_style:
        type: Optional[dict[str, Any]]
    """
    if table is not None and not isinstance(table, dict):
        raise TypeError("table must be a dictionary")
    if image is not None and not isinstance(image, (str, Path, dict)):
        raise TypeError("image must be a path string, Path, or dictionary")
    if markdown is not None and not isinstance(markdown, str):
        raise TypeError("markdown must be a string")
    if markdown and items:
        raise ValueError("markdown cannot be combined with items")
    if table and flow_boxes:
        raise ValueError("table cannot be combined with flow_boxes")
    if markdown and flow_boxes:
        raise ValueError("markdown cannot be combined with flow_boxes")
    if table and items and code:
        raise ValueError("table can be combined with items or code, not both")
    if table and markdown and (items or code):
        raise ValueError("table can be combined with one of items, markdown, or code")
    if image and flow_boxes:
        raise ValueError("image cannot be combined with flow_boxes")
    if image and table:
        raise ValueError("image cannot be combined with table")
    if image and items and code:
        raise ValueError("image can be combined with items or code, not both")
    if image and markdown and (items or code):
        raise ValueError("image can be combined with one of items, markdown, or code")

    slide_s = slide_style or {}
    code_defaults: dict[str, Any] = {
        "bg-color": "#193952",
        "font-color": "#FFFFFF",
        "font-size": 20,
        "font-name": "Consolas",
    }
    code_defaults.update(code_style or {})
    base_table_style = _merge_style(slide_s, table_style)

    section_gap = Inches(0.3)
    callout_height = Inches(0.8)
    reserved_callout = section_gap + callout_height if callout else 0
    main_height = max(Inches(1.0), CONTENT_HEIGHT - reserved_callout)
    bottom: int = CONTENT_TOP

    def place_image(image_top: int, image_height: int) -> None:
        assert image is not None

        fit = "contain"
        caption: str | None = None
        caption_style: Optional[dict[str, Any]] = None
        if isinstance(image, dict):
            image_path = _table_spec_value(image, "path", "src")
            fit_value = _table_spec_value(image, "fit", default="contain")
            caption = _table_spec_value(image, "caption")
            caption_style = _table_spec_value(
                image,
                "caption_style",
                "caption-style",
            )
            if not isinstance(fit_value, str):
                raise TypeError("image fit must be a string")
            fit = fit_value
        else:
            image_path = image

        if not isinstance(image_path, (str, Path)):
            raise TypeError("image path must be a string or Path")
        if caption is not None and not isinstance(caption, str):
            raise TypeError("image caption must be a string")
        if caption_style is not None and not isinstance(caption_style, dict):
            raise TypeError("image caption_style must be a dictionary")

        caption_gap = Inches(0.15)
        caption_height = Inches(0.6)
        reserved_caption = caption_gap + caption_height if caption else 0
        picture_height = max(1, image_height - reserved_caption)
        add_image(
            slide,
            left=CONTENT_LEFT,
            top=image_top,
            width=CONTENT_WIDTH,
            height=picture_height,
            image_file=image_path,
            fit=fit,
        )
        if caption:
            effective_caption_style = _merge_style(
                {
                    "font-size": 18,
                    "alignment": "center",
                    "italic": True,
                },
                slide_s,
            )
            effective_caption_style = _merge_style(
                effective_caption_style,
                caption_style,
            )
            add_textbox(
                slide,
                left=CONTENT_LEFT,
                top=image_top + picture_height + caption_gap,
                width=CONTENT_WIDTH,
                height=caption_height,
                text=caption,
                style=effective_caption_style,
            )

    def place_table(table_top: int, table_height: int) -> None:
        assert table is not None

        rows = _table_spec_value(table, "rows", default=[])
        if rows is None:
            rows = []
        if not isinstance(rows, list):
            raise TypeError("table rows must be a list of row lists")

        columns = _table_spec_value(table, "columns", "headers")
        if columns is not None and not isinstance(columns, list):
            raise TypeError("table columns must be a list")

        column_widths = _table_spec_value(table, "column_widths", "column-widths")
        if column_widths is not None and not isinstance(column_widths, list):
            raise TypeError("table column_widths must be a list")

        row_heights = _table_spec_value(table, "row_heights", "row-heights")
        if row_heights is not None and not isinstance(row_heights, list):
            raise TypeError("table row_heights must be a list")

        local_table_style = _table_spec_value(table, "style")
        if local_table_style is not None and not isinstance(local_table_style, dict):
            raise TypeError("table style must be a dictionary")
        local_header_style = _table_spec_value(table, "header_style", "header-style")
        if local_header_style is not None and not isinstance(local_header_style, dict):
            raise TypeError("table header_style must be a dictionary")
        local_cell_style = _table_spec_value(table, "cell_style", "cell-style")
        if local_cell_style is not None and not isinstance(local_cell_style, dict):
            raise TypeError("table cell_style must be a dictionary")

        effective_table_style = _merge_style(base_table_style, local_table_style)
        effective_header_style = _merge_style(
            _merge_style(effective_table_style, table_header_style),
            local_header_style,
        )
        effective_cell_style = _merge_style(
            _merge_style(effective_table_style, table_cell_style),
            local_cell_style,
        )

        add_table(
            slide,
            left=CONTENT_LEFT,
            top=table_top,
            width=CONTENT_WIDTH,
            height=table_height,
            rows=rows,
            columns=columns,
            column_widths=column_widths,
            row_heights=row_heights,
            banded_rows=_table_spec_value(table, "banded_rows", "banded-rows"),
            style=effective_table_style,
            header_style=effective_header_style,
            cell_style=effective_cell_style,
        )

    if flow_boxes:
        add_flow_boxes(
            slide,
            boxes=flow_boxes,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            style=slide_s,
        )
        bottom = CONTENT_TOP + Inches(3.0)

    elif markdown and code:
        markdown_height, code_height = _split_content_height(
            main_height,
            ratio=0.40,
            gap=section_gap,
            min_first=Inches(2.0),
            min_second=Inches(2.8),
        )
        add_markdown_textbox(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=markdown_height,
            markdown_text=markdown,
            style={**{"font-size": 24, "spacing": 12}, **slide_s},
        )
        code_top = CONTENT_TOP + markdown_height + section_gap
        add_code_block(
            slide,
            left=CONTENT_LEFT,
            top=code_top,
            width=CONTENT_WIDTH,
            height=code_height,
            code_text=code,
            style=code_defaults,
        )
        bottom = code_top + code_height

    elif items and code:
        bullets_height, code_height = _split_content_height(
            main_height,
            ratio=0.40,
            gap=section_gap,
            min_first=Inches(2.0),
            min_second=Inches(2.8),
        )
        add_bullet_list(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=bullets_height,
            items=items,
            style={**{"font-size": 26, "spacing": 10}, **slide_s},
        )
        code_top = CONTENT_TOP + bullets_height + section_gap
        add_code_block(
            slide,
            left=CONTENT_LEFT,
            top=code_top,
            width=CONTENT_WIDTH,
            height=code_height,
            code_text=code,
            style=code_defaults,
        )
        bottom = code_top + code_height

    elif items and table:
        bullets_height, table_height = _split_content_height(
            main_height,
            ratio=0.36,
            gap=section_gap,
            min_first=Inches(2.0),
            min_second=Inches(2.8),
        )
        add_bullet_list(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=bullets_height,
            items=items,
            style={**{"font-size": 26, "spacing": 10}, **slide_s},
        )
        table_top = CONTENT_TOP + bullets_height + section_gap
        place_table(table_top, table_height)
        bottom = table_top + table_height

    elif markdown and table:
        markdown_height, table_height = _split_content_height(
            main_height,
            ratio=0.36,
            gap=section_gap,
            min_first=Inches(2.0),
            min_second=Inches(2.8),
        )
        add_markdown_textbox(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=markdown_height,
            markdown_text=markdown,
            style={**{"font-size": 24, "spacing": 12}, **slide_s},
        )
        table_top = CONTENT_TOP + markdown_height + section_gap
        place_table(table_top, table_height)
        bottom = table_top + table_height

    elif code and table:
        code_height, table_height = _split_content_height(
            main_height,
            ratio=0.40,
            gap=section_gap,
            min_first=Inches(2.4),
            min_second=Inches(2.8),
        )
        add_code_block(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=code_height,
            code_text=code,
            style=code_defaults,
        )
        table_top = CONTENT_TOP + code_height + section_gap
        place_table(table_top, table_height)
        bottom = table_top + table_height

    elif items and image:
        bullets_height, image_height = _split_content_height(
            main_height,
            ratio=0.36,
            gap=section_gap,
            min_first=Inches(2.0),
            min_second=Inches(2.8),
        )
        add_bullet_list(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=bullets_height,
            items=items,
            style={**{"font-size": 26, "spacing": 10}, **slide_s},
        )
        image_top = CONTENT_TOP + bullets_height + section_gap
        place_image(image_top, image_height)
        bottom = image_top + image_height

    elif markdown and image:
        markdown_height, image_height = _split_content_height(
            main_height,
            ratio=0.36,
            gap=section_gap,
            min_first=Inches(2.0),
            min_second=Inches(2.8),
        )
        add_markdown_textbox(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=markdown_height,
            markdown_text=markdown,
            style={**{"font-size": 24, "spacing": 12}, **slide_s},
        )
        image_top = CONTENT_TOP + markdown_height + section_gap
        place_image(image_top, image_height)
        bottom = image_top + image_height

    elif code and image:
        code_height, image_height = _split_content_height(
            main_height,
            ratio=0.40,
            gap=section_gap,
            min_first=Inches(2.4),
            min_second=Inches(2.8),
        )
        add_code_block(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=code_height,
            code_text=code,
            style=code_defaults,
        )
        image_top = CONTENT_TOP + code_height + section_gap
        place_image(image_top, image_height)
        bottom = image_top + image_height

    elif table:
        place_table(CONTENT_TOP, main_height)
        bottom = CONTENT_TOP + main_height

    elif image:
        place_image(CONTENT_TOP, main_height)
        bottom = CONTENT_TOP + main_height

    elif markdown:
        add_markdown_textbox(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=main_height,
            markdown_text=markdown,
            style={**{"font-size": 26, "spacing": 14}, **slide_s},
        )
        bottom = CONTENT_TOP + main_height

    elif code:
        add_code_block(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=main_height,
            code_text=code,
            style=code_defaults,
        )
        bottom = CONTENT_TOP + main_height

    elif items:
        add_bullet_list(
            slide,
            left=CONTENT_LEFT,
            top=CONTENT_TOP,
            width=CONTENT_WIDTH,
            height=main_height,
            items=items,
            style={**{"font-size": 30, "spacing": 14}, **slide_s},
        )
        bottom = CONTENT_TOP + main_height

    if callout:
        callout_top = bottom + section_gap
        add_textbox(
            slide,
            left=CONTENT_LEFT,
            top=callout_top,
            width=CONTENT_WIDTH,
            height=callout_height,
            text=callout,
            style={**{"font-size": 26, "bold": True}, **slide_s},
        )


def clone_slide(prs: Presentation, template_idx: int) -> Slide:
    """
    title: Clone a slide from the presentation by index.
    summary: |-
      Creates a deep copy of the slide's XML and all
      relationships, appending the new slide at the end.
    parameters:
      prs:
        type: Presentation
        description: The presentation object.
      template_idx:
        type: int
        description: Zero-based index of the slide to clone.
    returns:
      type: Slide
      description: The newly created slide.
    """
    template_slide = prs.slides[template_idx]
    slide_layout = template_slide.slide_layout

    new_slide = prs.slides.add_slide(slide_layout)

    # Remove default shapes from the new slide
    for shape in list(new_slide.shapes):
        sp_elem = shape._element
        sp_elem.getparent().remove(sp_elem)

    # Copy all elements from template slide
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # Copy slide background
    if template_slide.background._element is not None:
        # Use the csld element's bg
        csld = new_slide._element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}cSld"
        )
        if csld is not None:
            old_bg = csld.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}bg"
            )
            if old_bg is not None:
                csld.remove(old_bg)

    # Copy relationships (images, etc.)
    for rel in template_slide.part.rels.values():
        if "image" in rel.reltype:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)

    return new_slide


def delete_slide(prs: Presentation, slide_idx: int) -> None:
    """
    title: Delete a slide from the presentation by index.
    parameters:
      prs:
        type: Presentation
        description: The presentation object.
      slide_idx:
        type: int
        description: Zero-based index of the slide to delete.
    """
    rId = prs.slides._sldIdLst[slide_idx].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[slide_idx]
    prs.slides._sldIdLst.remove(sldId)


def move_slide(prs: Presentation, old_idx: int, new_idx: int) -> None:
    """
    title: Move a slide from one position to another.
    parameters:
      prs:
        type: Presentation
        description: The presentation object.
      old_idx:
        type: int
        description: Current zero-based index of the slide.
      new_idx:
        type: int
        description: Desired zero-based index for the slide.
    """
    sld_id_lst = prs.slides._sldIdLst
    el = sld_id_lst[old_idx]
    sld_id_lst.remove(el)
    if new_idx >= len(sld_id_lst):
        sld_id_lst.append(el)
    else:
        sld_id_lst.insert(new_idx, el)
