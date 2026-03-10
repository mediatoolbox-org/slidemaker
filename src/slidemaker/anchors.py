"""Anchor-map helpers for template-specific shape and layout bindings."""

from __future__ import annotations

import json
from copy import deepcopy
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation

DEFAULT_TEMPLATE_PAGE = 5


def default_anchor_map(
    default_template_page: int = DEFAULT_TEMPLATE_PAGE,
) -> dict[str, Any]:
    """Return the built-in anchor map matching the default template."""
    return {
        "version": 1,
        "default-template-page": int(default_template_page),
        "title-slide": {
            "title-group": "Group 7",
            "subtitle-group": "Group 4",
        },
        "title-groups": {
            1: "Group 7",
            2: "Group 3",
            3: "Group 4",
            4: "Group 4",
            5: "Group 4",
            6: "Group 4",
            7: "Group 4",
            8: "Group 4",
            9: "Group 4",
        },
        "remove-shapes": {
            "objectives": ["TextBox 6"],
        },
        "checkpoints": {
            "textbox-names": [
                "TextBox 34",
                "TextBox 35",
                "TextBox 36",
                "TextBox 37",
                "TextBox 38",
            ],
        },
        "areas": {
            "objectives-bullets": {
                "left": 0.9,
                "top": 3.6,
                "width": 13.0,
                "height": 6.5,
            },
            "toolkit-bullets": {
                "left": 0.9,
                "top": 2.6,
                "width": 12.0,
                "height": 7.5,
            },
            "whats-new-bullets": {
                "left": 0.9,
                "top": 2.6,
                "width": 12.0,
                "height": 7.5,
            },
            "generic-flow": {
                "left": 0.9,
                "top": 2.6,
            },
            "generic-bullets-code-bullets": {
                "left": 0.9,
                "top": 2.6,
                "width": 8.5,
                "height": 3.5,
            },
            "generic-bullets-code-code": {
                "left": 0.9,
                "top": 6.4,
                "width": 13.0,
                "height": 4.5,
            },
            "generic-code-only": {
                "left": 0.9,
                "top": 2.6,
                "width": 13.0,
                "height": 6.0,
            },
            "generic-bullets-only": {
                "left": 0.9,
                "top": 2.6,
                "width": 12.5,
                "height": 7.5,
            },
            "generic-callout": {
                "left": 0.9,
                "width": 12.5,
                "height": 0.8,
                "offset-top": 0.3,
            },
            "exercise-bullets": {
                "left": 0.9,
                "top": 2.6,
                "width": 12.0,
                "height": 7.5,
            },
            "debugging-bullets": {
                "left": 0.9,
                "top": 2.6,
                "width": 12.0,
                "height": 7.5,
            },
            "recap-bullets": {
                "left": 0.9,
                "top": 2.6,
                "width": 12.0,
                "height": 5.0,
            },
            "recap-next-topic": {
                "left": 0.9,
                "top": 8.5,
                "width": 12.0,
                "height": 1.5,
            },
        },
    }


def _shape_catalog(template: str | Path) -> dict[int, list[dict[str, Any]]]:
    """Inspect template slides and return a simple shape catalog."""
    prs = Presentation(str(template))
    catalog: dict[int, list[dict[str, Any]]] = {}
    for page, slide in enumerate(prs.slides, start=1):
        entries: list[dict[str, Any]] = []
        for shape in slide.shapes:
            item: dict[str, Any] = {
                "name": shape.name,
                "shape-type": int(shape.shape_type),
                "has-text-frame": bool(getattr(shape, "has_text_frame", False)),
            }
            if item["has-text-frame"]:
                preview = shape.text_frame.text.strip().replace("\n", " ")
                if preview:
                    item["text-preview"] = preview[:120]
            if shape.shape_type == 6:  # GROUP
                children = [child.name for child in shape.shapes]
                if children:
                    item["children"] = children
            entries.append(item)
        catalog[page] = entries
    return catalog


def generate_anchor_map(
    template: str | Path,
    default_template_page: int = DEFAULT_TEMPLATE_PAGE,
    include_shape_catalog: bool = True,
) -> dict[str, Any]:
    """Generate an editable anchor map initialized from built-in defaults."""
    anchor_map = default_anchor_map(default_template_page)
    anchor_map["template"] = str(template)
    if include_shape_catalog:
        anchor_map["shape-catalog"] = _shape_catalog(template)
    return anchor_map


def _load_yaml(text: str) -> dict[str, Any]:
    """Load YAML text using PyYAML when available."""
    try:
        import yaml  # type: ignore
    except ModuleNotFoundError as exc:
        raise RuntimeError(
            "YAML parsing requires PyYAML. Install it with "
            "`pip install pyyaml` or use JSON content in the file."
        ) from exc
    loaded = yaml.safe_load(text)
    if not isinstance(loaded, dict):
        raise TypeError("anchor map file must contain a mapping/object at root")
    return loaded


def load_anchor_map(
    anchor_map: Optional[dict[str, Any] | str | Path],
) -> Optional[dict[str, Any]]:
    """Load anchor map from dict or file path."""
    if anchor_map is None:
        return None
    if isinstance(anchor_map, dict):
        return deepcopy(anchor_map)
    path = Path(anchor_map)
    text = path.read_text(encoding="utf-8")
    try:
        loaded = json.loads(text)
        if not isinstance(loaded, dict):
            raise TypeError("anchor map JSON must be an object")
        return loaded
    except json.JSONDecodeError:
        return _load_yaml(text)


def dump_anchor_map(anchor_map: dict[str, Any], out: str | Path) -> Path:
    """Write anchor map to disk as YAML if possible, else JSON-compatible YAML."""
    out_path = Path(out)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        import yaml  # type: ignore

        rendered = yaml.safe_dump(
            anchor_map,
            sort_keys=False,
            allow_unicode=False,
            width=100,
        )
    except ModuleNotFoundError:
        # JSON is valid YAML 1.2 and keeps this feature dependency-free.
        rendered = json.dumps(anchor_map, indent=2)
    out_path.write_text(rendered, encoding="utf-8")
    return out_path
