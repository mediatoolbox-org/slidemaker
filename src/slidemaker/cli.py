"""
title: High-level API for building slide decks.
summary: |-
  ``SlideBuilder`` is the main entry point.  A user script creates
  an instance, calls ``add_slide`` to clone and populate slides, and
  finally calls :meth:`save` to write the ``.pptx`` file.

  Example
  -------
  ::

  from slidemaker import SlideBuilder

  sb = SlideBuilder("template.pptx")
  sb.add_slide(content={"title": "Hello", "body": "World"})
  sb.add_slide(items=["Extract", "Transform", "Load"])
  sb.save("output.pptx")
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation

from slidemaker.core import (
    clone_slide,
    delete_slide,
    layout_content_shapes,
    remove_generated_content_placeholders,
    replace_placeholders,
    set_notes,
)

DEFAULT_TEMPLATE_PAGE = 5
_SYSTEM_STYLES = {".slide", ".code", ".table", ".table-header", ".table-cell"}

StyleAttrs = dict[str, Any]
StyleMap = dict[str, StyleAttrs]


class SlideBuilder:
    """
    title: Builds a slide deck from a PowerPoint template.
    summary: |-
      Each ``add_slide`` call clones a template slide and populates it
      with content.  Call :meth:`save` to assemble the final deck.

      **Style keys**

      - ``.slide`` — base style for new text shapes (bullets, callout).
      - ``.code``  — style for code blocks.
      - ``.table`` — base style for generated tables.
      - ``.table-header`` — header-row overrides for generated tables.
      - ``.table-cell`` — body-cell overrides for generated tables.
      - ``#placeholder`` — style applied when replacing a
        ``{{placeholder}}`` in the template.  Falls back to ``.slide``.

      Named styles (registered via ``add_style``) can be referenced
      by string name in ``add_slide(style="dense")``.
    attributes:
      _template_path:
        description: Path to the template file.
      _prs:
        description: The python-pptx Presentation being built.
      _template_default_page:
        description: 1-based index of the template slide to clone by default.
      _template_slide_count:
        description: Number of slides in the original template.
      _slide_count:
        description: Number of content slides added so far.
      _styles:
        type: StyleMap
        description: Registered styles (system and placeholder).
    """

    def __init__(
        self,
        template: str | Path,
        style: Optional[StyleMap] = None,
        template_default_page: int = DEFAULT_TEMPLATE_PAGE,
    ) -> None:
        """
        title: Load template, styles, and prepare for building.
        parameters:
          template:
            type: str | Path
          style:
            type: Optional[StyleMap]
          template_default_page:
            type: int
        """
        self._template_path = Path(template)
        self._prs = Presentation(str(self._template_path))
        self._template_default_page = template_default_page
        self._template_slide_count = len(self._prs.slides)
        self._slide_count = 0
        self._styles: StyleMap = {
            ".slide": {},
            ".code": {},
            ".table": {},
            ".table-header": {},
            ".table-cell": {},
        }
        if style:
            self.add_style(style)

    def add_style(self, style: StyleMap) -> None:
        """
        title: Register or update style definitions.
        summary: |-
          Accepted key prefixes:

          - ``.slide`` — base text style.
          - ``.code``  — code block style.
          - ``.table`` — base generated-table style.
          - ``.table-header`` — generated-table header style.
          - ``.table-cell`` — generated-table body-cell style.
          - ``#name``  — placeholder style for ``{{name}}``.
          - Any other name — a named style preset that can be
            referenced by string in ``add_slide(style="name")``.
        parameters:
          style:
            type: StyleMap
            description: >-
              Mapping of style names to style attribute dictionaries.
        """
        for name, attrs in style.items():
            if not isinstance(name, str):
                raise TypeError("style name keys must be strings")
            if not isinstance(attrs, dict):
                raise TypeError(
                    f"style '{name}' must map to a dictionary of attributes"
                )
            if name not in self._styles:
                self._styles[name] = {}
            self._styles[name].update(attrs)

    def _resolve_styles(
        self,
        style: Optional[str | dict[str, Any]],
    ) -> StyleMap:
        """
        title: Resolve effective styles for one slide call.
        summary: |-
          Returns a merged StyleMap containing ``.slide``, ``.code``,
          generated-table styles, and any ``#placeholder`` overrides.

          - ``None``  → global styles only.
          - ``str``   → apply a named preset to ``.slide``.
          - ``dict``  → merge per-slide overrides by key prefix.
        parameters:
          style:
            type: Optional[str | dict[str, Any]]
        returns:
          type: StyleMap
        """
        styles: StyleMap = deepcopy(self._styles)

        if style is None:
            return styles

        if isinstance(style, str):
            # Named preset — merge into .slide
            if style not in self._styles:
                raise KeyError(f"unknown style name: {style}")
            styles[".slide"].update(self._styles[style])
            return styles

        if not isinstance(style, dict):
            raise TypeError("style must be None, a style name, or a dict")

        # Check if any values are dicts (structured style override)
        has_nested = any(isinstance(v, dict) for v in style.values())

        if has_nested:
            # Apply named preset(s) via "use" key
            use_names = style.get("use")
            if isinstance(use_names, str):
                if use_names not in self._styles:
                    raise KeyError(f"unknown style name: {use_names}")
                styles[".slide"].update(self._styles[use_names])
            elif isinstance(use_names, list):
                for name in use_names:
                    if isinstance(name, str) and name in self._styles:
                        styles[".slide"].update(self._styles[name])

            for key, attrs in style.items():
                if key == "use" or not isinstance(key, str):
                    continue
                if isinstance(attrs, dict):
                    # .slide, .code, #placeholder keys
                    styles.setdefault(key, {}).update(attrs)
                else:
                    # Bare key-value → goes into .slide
                    styles[".slide"][key] = attrs
            return styles

        # Flat dict → all goes into .slide
        styles[".slide"].update(style)
        return styles

    # ── Public API ────────────────────────────────────────────

    def add_slide(
        self,
        content: dict[str, str | list[str] | None] | None = None,
        items: list[str] | None = None,
        markdown: str | None = None,
        code: str | None = None,
        table: dict[str, Any] | None = None,
        image: str | Path | dict[str, Any] | None = None,
        flow_boxes: list[dict] | None = None,
        callout: str | None = None,
        notes: str = "",
        style: Optional[str | dict[str, Any]] = None,
        template_page: int | None = None,
    ) -> None:
        """
        title: Add a content slide by cloning a template page.
        summary: |-
          Two content modes can be used independently or together:

          - **Replace** — pass ``content`` to replace existing
            ``{{placeholder}}`` text in the cloned template shapes.
          - **Create** — pass ``items``, ``markdown``, ``code``, ``table``,
            ``image``, ``flow_boxes``, or ``callout`` to add new shapes on top
            of the cloned slide.
        parameters:
          content:
            type: dict[str, str | list[str] | None] | None
            description: >-
              Mapping of placeholder names to replacement values. Shapes
              containing ``{{key}}`` are replaced with the corresponding value
              (str, list of bullets, or None to clear).
          items:
            type: list[str] | None
            description: Bullet point strings (creates new shapes).
          markdown:
            type: str | None
            description: Markdown text block (creates a new shape).
          code:
            type: str | None
            description: Source code for a dark code block.
          table:
            type: dict[str, Any] | None
            description: >-
              Table specification for a generated table shape. Supports
              ``columns`` (optional header row), ``rows``, optional
              ``column_widths``/``row_heights``, optional ``banded_rows``, and
              optional ``style``/``header_style``/``cell_style``.
          image:
            type: str | Path | dict[str, Any] | None
            description: >-
              Image path or image specification for a generated picture shape.
              Dict form supports ``path``/``src``, optional ``caption``,
              optional ``fit`` (``contain`` or ``stretch``), and optional
              ``caption_style``.
          flow_boxes:
            type: list[dict] | None
            description: Flow-diagram boxes (label, desc, optional style).
          callout:
            type: str | None
            description: Bold callout line below other content.
          notes:
            type: str
            description: Speaker notes.
          style:
            type: Optional[str | dict[str, Any]]
            description: >-
              Per-slide style override. Supports ``.slide``, ``.code``,
              ``.table``, ``.table-header``, ``.table-cell``, and
              ``#placeholder`` keys.
          template_page:
            type: int | None
            description: >-
              1-based template page to clone. Defaults to
              ``template_default_page`` from the constructor.
        """
        styles = self._resolve_styles(style)
        page = (
            template_page if template_page is not None else self._template_default_page
        )
        idx = page - 1
        new_slide = clone_slide(self._prs, idx)
        self._slide_count += 1

        # ── Replace placeholders ──────────────────────────────
        if content is not None:
            replace_placeholders(new_slide, content, styles=styles)

        # ── Create new shapes with smart layout ───────────────
        if items or markdown or code or table or image or flow_boxes or callout:
            remove_generated_content_placeholders(new_slide)
            layout_content_shapes(
                new_slide,
                items=items,
                markdown=markdown,
                code=code,
                table=table,
                image=image,
                flow_boxes=flow_boxes,
                callout=callout,
                slide_style=styles.get(".slide"),
                code_style=styles.get(".code"),
                table_style=styles.get(".table"),
                table_header_style=styles.get(".table-header"),
                table_cell_style=styles.get(".table-cell"),
            )

        if notes:
            set_notes(new_slide, notes)

    # ── Build and save ──────────────────────────────────────

    def save(self, path: str) -> None:
        """
        title: Assemble the final deck and write to disk.
        summary: Removes all original template slides and saves.
        parameters:
          path:
            type: str
            description: Output file path for the .pptx file.
        """
        for idx in range(self._template_slide_count - 1, -1, -1):
            delete_slide(self._prs, idx)

        self._prs.save(path)
        n_slides = len(self._prs.slides)
        print(f"Saved → {path}  ({n_slides} slides)")
